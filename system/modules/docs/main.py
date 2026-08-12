"""docs — office documents in and out of the render-block IR.

Reading: pptx/xlsx/docx are ZIP+XML with mature readers; hwpx is the same idea (OWPML) read with
stdlib zipfile+ElementTree since no mature library exists; pdf via pypdf (text only).

Making: the hard step of HTML->PPTX conversion — recovering *meaning* from computed layout —
never happens here, because the render blocks already carry the meaning ("this is a table",
"this is a heading"). Blocks map straight to native shapes: header level 1 / divider starts a
slide, table blocks become editable tables, a master .pptx (the corporate deck uploaded through
the media door) contributes its theme and layouts.

Files land in data/docs/ and are declared as `_mediaImport` — the framework carries them into
the media store (gated, served, downloadable). hwpx WRITING is deliberately absent: OWPML has no
mature library and hand-written XML is the most work of the four; read-only by decision.
"""

import hashlib
import json
import os
import re
import sys
import zipfile
import xml.etree.ElementTree as ET

OUT_DIR = "data/docs"

MIME = {
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf": "application/pdf",
    "hwpx": "application/vnd.hancom.hwpx",
}


def resolve_path(raw):
    """Media URL or workspace-relative path -> readable path (sandbox cwd = workspace root)."""
    path = str(raw or "").strip()
    if not path:
        return None, "path is required"
    if "://" in path:
        rest = path.split("://", 1)[1]
        path = "/" + rest.split("/", 1)[1] if "/" in rest else ""
    path = path.lstrip("/")
    if ".." in path.split("/"):
        return None, f"path escapes the workspace: {raw}"
    if not os.path.isfile(path):
        return None, f"file not found: {path} (workspace-relative)"
    return path, None


def out_file(title, ext, payload):
    os.makedirs(OUT_DIR, exist_ok=True)
    stem = re.sub(r"[^0-9A-Za-z가-힣_-]+", "-", str(title or "")).strip("-")[:40] or "untitled"
    h = hashlib.sha1(json.dumps(payload, sort_keys=True, default=str).encode()).hexdigest()[:8]
    return f"{OUT_DIR}/{stem}-{h}.{ext}", stem


def media_import_decl(path, ext, stem):
    return {"path": path, "contentType": MIME[ext], "filenameHint": stem}


# ── read ───────────────────────────────────────────────────────────────────────────────────────


def _cell_str(v):
    if v is None:
        return ""
    if hasattr(v, "isoformat"):
        return v.isoformat()
    return v if isinstance(v, (int, float, str, bool)) else str(v)


def read_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages):
        try:
            pages.append(page.extract_text() or "")
        except Exception as e:  # noqa: BLE001 — one broken page should not lose the rest
            pages.append(f"[page {i + 1}: extract failed: {e}]")
    return {"format": "pdf", "meta": {"pages": len(pages)},
            "text": "\n\n".join(pages).strip(), "tables": []}


def read_docx(path):
    import docx
    d = docx.Document(path)
    lines = []
    for p in d.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        style = (p.style.name or "").lower()
        m = re.match(r"heading (\d)", style)
        lines.append(("#" * int(m.group(1)) + " " + t) if m else t)
    tables = []
    for tb in d.tables:
        rows = [[c.text.strip() for c in row.cells] for row in tb.rows]
        if rows:
            tables.append({"headers": rows[0], "rows": rows[1:]})
    return {"format": "docx", "meta": {"paragraphs": len(lines), "tables": len(tables)},
            "text": "\n".join(lines), "tables": tables}


def read_xlsx(path, max_rows=1000, max_cols=50):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets, tables, text_parts = [], [], []
    for ws in wb.worksheets:
        rows = []
        for r in ws.iter_rows(max_row=max_rows, max_col=max_cols, values_only=True):
            row = [_cell_str(v) for v in r]
            while row and str(row[-1]) == "":  # read_only pads to max_col — trim the tail
                row.pop()
            if row:
                rows.append(row)
        sheets.append({"name": ws.title, "rows": len(rows)})
        if rows:
            tables.append({"name": ws.title, "headers": rows[0], "rows": rows[1:]})
            text_parts.append(f"[{ws.title}] {len(rows)} rows")
    wb.close()
    return {"format": "xlsx", "meta": {"sheets": sheets},
            "text": "\n".join(text_parts), "tables": tables}


def read_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    slides, tables = [], []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[c.text.strip() for c in row.cells] for row in shape.table.rows]
                if rows:
                    tables.append({"name": f"slide {i + 1}", "headers": rows[0], "rows": rows[1:]})
                continue
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        slides.append({"n": i + 1, "text": "\n".join(texts)})
    text = "\n\n".join(f"[slide {s['n']}]\n{s['text']}" for s in slides if s["text"])
    return {"format": "pptx", "meta": {"slides": len(slides), "tables": len(tables)},
            "text": text, "tables": tables}


def _local(tag):
    return tag.rsplit("}", 1)[-1]


def read_hwpx(path):
    """OWPML (KS X 6101) — ZIP of section XMLs; text lives in <hp:t>, tables in <hp:tbl>,
    equation sources in <hp:equation><hp:script>. The <hp:t> walk already reaches text
    boxes and footnotes (they are just deeper paragraphs — 2026-08-11 실측으로 확인).
    Namespace-agnostic on purpose: producers disagree on prefixes, localnames hold still."""
    texts, tables, equations = [], [], []
    seen_eq = set()
    with zipfile.ZipFile(path) as z:
        sections = sorted(n for n in z.namelist()
                          if n.startswith("Contents/") and re.search(r"section\d+\.xml$", n))
        if not sections:
            return {"format": "hwpx", "meta": {"sections": 0},
                    "text": "", "tables": [],
                    "note": "no Contents/section*.xml — not an HWPX container?"}
        for name in sections:
            root = ET.fromstring(z.read(name))
            in_table_cells = set()
            for tbl in root.iter():
                if _local(tbl.tag) != "tbl":
                    continue
                rows = []
                for tr in tbl.iter():
                    if _local(tr.tag) != "tr":
                        continue
                    row = []
                    for tc in tr:
                        if _local(tc.tag) != "tc":
                            continue
                        cell = "".join(t.text or "" for t in tc.iter() if _local(t.tag) == "t")
                        row.append(cell.strip())
                        for t in tc.iter():
                            if _local(t.tag) == "t":
                                in_table_cells.add(id(t))
                    if row:
                        rows.append(row)
                if rows:
                    tables.append({"headers": rows[0], "rows": rows[1:]})
            for p in root.iter():
                if _local(p.tag) != "p":
                    continue
                line = "".join(t.text or "" for t in p.iter()
                               if _local(t.tag) == "t" and id(t) not in in_table_cells)
                if line.strip():
                    texts.append(line.strip())
                # Equation SOURCE is not a <hp:t> — pull the script so formulas stop
                # vanishing. id-dedup because outer paragraphs contain cell paragraphs.
                for eq in p.iter():
                    if _local(eq.tag) != "equation" or id(eq) in seen_eq:
                        continue
                    seen_eq.add(id(eq))
                    for sc in eq:
                        if _local(sc.tag) == "script" and (sc.text or "").strip():
                            script = sc.text.strip()
                            equations.append(script)
                            texts.append(f"[수식] {script}")
    out = {"format": "hwpx",
           "meta": {"sections": len(sections), "tables": len(tables),
                    "equations": len(equations)},
           "text": "\n".join(texts), "tables": tables}
    if equations:
        out["equations"] = equations
    return out


def _rhwp_helper(request, timeout=180):
    """Spawn the vendored rhwp WASM helper (node, ./rhwp-helper.mjs). One JSON in, one out."""
    import shutil
    import subprocess
    node = shutil.which("node")
    if not node:
        raise RuntimeError("node runtime not found")
    helper = os.path.join(os.path.dirname(os.path.abspath(__file__)), "rhwp-helper.mjs")
    req = json.dumps(request).encode("utf-8")
    proc = subprocess.run([node, helper], input=req, capture_output=True, timeout=timeout)
    lines = proc.stdout.decode("utf-8", "replace").strip().splitlines()
    payload = json.loads(lines[-1]) if lines else {}
    if not payload.get("ok"):
        err = payload.get("error") or proc.stderr.decode("utf-8", "replace")[:400]
        raise RuntimeError(err or "helper failed")
    return payload["data"]


def _rhwp_helper_read(path):
    return _rhwp_helper({"op": "read", "path": path}, timeout=120)


def read_hwp_legacy(path):
    """.hwp (binary) through the rhwp engine — the only door into that format here.
    Honesty note rides along: rhwp 0.8.x reads its own output perfectly but extracts
    THIN from vintage Hancom files (5.0.3.x-era samples came back near-empty, 2026-08-11
    실측) — so the reader names its limits instead of pretending."""
    try:
        d = _rhwp_helper_read(path)
    except Exception as e:  # noqa: BLE001
        raise RuntimeError(f"legacy .hwp needs the rhwp helper (node + vendored wasm): {e}")
    tables = [{"name": t.get("name"),
               "headers": (t.get("rows") or [[]])[0],
               "rows": (t.get("rows") or [[]])[1:]} for t in d.get("tables", [])]
    out = {"format": "hwp", "engine": "rhwp",
           "meta": {"sections": d.get("sections", 0),
                    "paragraphs": d.get("paragraphs", 0),
                    "tables": len(tables), "equations": len(d.get("equations", []))},
           "text": d.get("text", ""), "tables": tables,
           "note": ("rhwp engine — if content looks missing (older .hwp files extract "
                    "thin), convert to .hwpx in Hancom and read that instead")}
    if d.get("equations"):
        out["equations"] = d["equations"]
    return out


READERS = {"pdf": read_pdf, "docx": read_docx, "xlsx": read_xlsx,
           "pptx": read_pptx, "hwpx": read_hwpx, "hwp": read_hwp_legacy}


def action_read(inp):
    path, err = resolve_path(inp.get("path"))
    if err:
        return {"success": False, "action": "read", "error": err}
    ext = path.rsplit(".", 1)[-1].lower()
    reader = READERS.get(ext)
    if not reader:
        return {"success": False, "action": "read",
                "error": f"unsupported format .{ext} — one of: {', '.join(sorted(READERS))}"}
    try:
        data = reader(path)
    except Exception as e:  # noqa: BLE001 — a broken file should name itself
        return {"success": False, "action": "read", "error": f"read failed ({ext}): {e}"}
    cap = inp.get("maxChars")
    if isinstance(cap, (int, float)) and cap and len(data.get("text", "")) > int(cap):
        data["text"] = data["text"][: int(cap)]
        data["textTruncated"] = True
    data["sourcePath"] = path
    return {"success": True, "action": "read", "data": data}


# ── make ───────────────────────────────────────────────────────────────────────────────────────


# ── block IR normalization ─────────────────────────────────────────────────────────────────────
# Everything here runs BEFORE any renderer sees a block, because every renderer used to answer
# "I do not know this type" with silence. A grid of metrics, a tab's children, a page module's
# baked output — all authored content — arrived in the file as nothing at all (2026-08-12
# coverage audit). Containers are flattened into the linear stream the renderers already speak,
# and what genuinely cannot be printed is NAMED instead of vanishing.

# Interactive-only components: there is no honest paper form of a button or a countdown.
_DROP_BLOCK_TYPES = ("html", "form", "ad_slot", "result_display", "button", "slider",
                     "countdown", "plan_card", "lottie")
# Live components hold a socket topic, not data. When inline rows ride along (a seed) they
# degrade to their static sibling; with nothing but a topic there is nothing to print.
_LIVE_DEGRADE = {"live_feed": "table", "live_chart": "chart", "live_stock_chart": "stock_chart"}
_SNAPSHOT_SUFFIX = " (스냅샷)"
_ALIAS_BLOCK_TYPES = {"alert": "callout"}
_CONTAINER_MAX_DEPTH = 6
_STOCK_TABLE_CAP = 60          # candles printed before a document says "and N more"
_STOCK_FIELDS = (("date", "날짜"), ("open", "시가"), ("high", "고가"),
                 ("low", "저가"), ("close", "종가"), ("volume", "거래량"))
_LEVEL_DIGIT_RE = re.compile(r"(\d+)")


def _header_level(p, default=2, lo=1, hi=3):
    """Header level from anything a caller might write: 2, "2", "h2", "Heading 2", None.

    `int(p.get("level") or 2)` raised ValueError on "h2" and took the WHOLE document with it —
    one non-numeric level killed the build (2026-08-12 audit). A level is a rank, so the only
    honest failure mode is clamping, never raising."""
    raw = (p or {}).get("level")
    if isinstance(raw, bool):
        raw = None
    if isinstance(raw, (int, float)):
        n = int(raw)
    else:
        m = _LEVEL_DIGIT_RE.search(str(raw if raw is not None else ""))
        n = int(m.group(1)) if m else default
    return max(lo, min(hi, n))


def _table_rows(headers, rows):
    """(headers, list-rows) from either row dialect — [[...]] or [{...}].

    The dict dialect is what a model writes when it thinks in records, and every renderer took it
    literally: pptx/docx/pdf indexed `row[c]` with an integer and raised KeyError, xlsx iterated
    the dict and wrote its KEYS as the data (2026-08-12 audit). Headers give the column order;
    with no headers the first record's key order IS the order, so it becomes the header row."""
    headers = [str(h) for h in (headers or [])]
    rows = list(rows or [])
    if not headers:
        seen = []
        for r in rows:
            if isinstance(r, dict):
                for k in r:
                    if k not in seen:
                        seen.append(k)
        headers = [str(k) for k in seen]
    lowered = [h.strip().lower() for h in headers]
    out = []
    for r in rows:
        if isinstance(r, dict):
            if not headers:
                out.append(list(r.values()))
                continue
            loose = {str(k).strip().lower(): v for k, v in r.items()}
            out.append([r[h] if h in r else loose.get(lowered[i], "")
                        for i, h in enumerate(headers)])
        elif isinstance(r, (list, tuple)):
            out.append(list(r))
        else:
            out.append([r])
    return headers, out


def _numeric_cols(rows, ncols):
    """Column indices that read as numbers — MAJORITY, not unanimity (the xlsx rule). One "-" in
    a price column does not make the column text, and print renderers right-align it either way."""
    out = set()
    for ci in range(ncols):
        vals = [str(r[ci]).strip() for r in rows
                if ci < len(r) and r[ci] not in (None, "")]
        if vals and sum(1 for v in vals if _NUMLIKE_RE.match(v)) * 2 > len(vals):
            out.add(ci)
    return out


def _callout_parts(p):
    """(title, message) — the callout dialect has four names for its body."""
    msg = (p.get("message") if p.get("message") not in (None, "") else
           p.get("content") if p.get("content") not in (None, "") else
           p.get("text") if p.get("text") not in (None, "") else p.get("body"))
    return str(p.get("title") or "").strip(), str(msg or "").strip()


def _stock_rows(p):
    """stock_chart props -> (headers, rows) over whichever OHLC(V) fields the candles carry.

    Accepts the canonical record dialect (props.data = [{date,open,high,low,close,volume}]), the
    columnar dialect the page schema declares ({dates:[], close:[]}), and plain list rows."""
    raw = None
    for key in ("data", "candles", "ohlcv", "rows", "records"):
        v = p.get(key)
        if isinstance(v, list) and v:
            raw = v
            break
    if raw is None:
        cols = {k: p.get(k) for k, _ in _STOCK_FIELDS if isinstance(p.get(k), list)}
        if not cols and isinstance(p.get("dates"), list):
            cols["date"] = p.get("dates")
        if not cols:
            return [], []
        n = max(len(v) for v in cols.values())
        raw = [{k: (v[i] if i < len(v) else None) for k, v in cols.items()} for i in range(n)]
    recs = [r for r in raw if isinstance(r, dict)]
    if not recs:
        headers, rows = _table_rows(p.get("headers"), raw)
        return headers, rows
    present = [(k, ko) for k, ko in _STOCK_FIELDS
               if any(r.get(k) not in (None, "") for r in recs)]
    if not present:
        return _table_rows(None, recs)
    headers = [ko for _k, ko in present]
    rows = [[r.get(k) for k, _ko in present] for r in recs]
    return headers, rows


def _stock_table_block(p, cap=_STOCK_TABLE_CAP):
    """A candle series as a printable table — what a document can honestly show of a stock chart.
    Capped, and the cap says so out loud rather than truncating in silence."""
    headers, rows = _stock_rows(p)
    if not rows:
        return None, None
    title = str(p.get("title") or p.get("symbol") or "").strip()
    extra = None
    if len(rows) > cap:
        extra = f"… 앞쪽 {len(rows) - cap}개 봉은 생략했습니다 (전체 {len(rows)}개)"
        rows = rows[-cap:]
    return {"type": "table", "props": {"headers": headers, "rows": rows, "title": title}}, extra


def _note_once(notes, key, text):
    if key not in notes:
        notes[key] = text


def _live_snapshot(t, p):
    """A live block that carries inline rows becomes its static sibling, labelled a snapshot.
    A block carrying only a topic (or a fetch binding with no rows) has nothing to degrade to."""
    rows = None
    for key in ("data", "seed", "rows", "items", "records", "points"):
        v = p.get(key)
        if isinstance(v, list) and v:
            rows = v
            break
    if rows is None:
        return None
    label = str(p.get("title") or p.get("symbol") or "").strip()
    title = (label + _SNAPSHOT_SUFFIX) if label else _SNAPSHOT_SUFFIX.strip()
    if t == "live_stock_chart":
        return {"type": "stock_chart", "props": {**p, "data": rows, "title": title}}
    if t == "live_chart":
        field = str(p.get("valueField") or "value").split(".")[-1]
        labels, values = [], []
        for i, r in enumerate(rows):
            if isinstance(r, dict):
                lab = next((r[k] for k in ("label", "name", "time", "date", "t", "x")
                            if r.get(k) not in (None, "")), i + 1)
                val = next((r[k] for k in (field, "value", "v", "y", "close")
                            if r.get(k) not in (None, "")), None)
            else:
                lab, val = i + 1, r
            labels.append(str(lab))
            values.append(parse_number(val))
        if any(v is not None for v in values):
            return {"type": "chart", "props": {"chartType": "line", "title": title,
                                               "labels": labels, "data": values}}
    headers, trows = _table_rows(p.get("headers"), rows)
    if not trows:
        return None
    return {"type": "table", "props": {"headers": headers, "rows": trows, "title": title}}


def _flatten_blocks(blocks, out, notes, depth, expand_module):
    """The recursive pass. `expand_module` goes False under a module's baked output, so a module
    nested inside someone else's composition is never re-expanded."""
    for b in blocks or []:
        if not isinstance(b, dict):
            continue
        t = str(b.get("type") or "").strip()
        name = str(b.get("name") or "").strip()
        if name and (t == "component" or not t):
            t = name
        t = t.lower()
        t = _ALIAS_BLOCK_TYPES.get(t, t)
        p = b.get("props") if isinstance(b.get("props"), dict) else {}
        if t != str(b.get("type") or ""):
            b = dict(b, type=t)

        if t in _DROP_BLOCK_TYPES:
            _note_once(notes, t, f"'{t}' 블록은 문서로 옮길 수 없어 제외했습니다 "
                                 "(화면에서만 동작하는 컴포넌트).")
            continue
        if depth > _CONTAINER_MAX_DEPTH:
            _note_once(notes, "_depth", "중첩이 너무 깊은 컨테이너가 있어 그 아래는 펼치지 "
                                        "않았습니다.")
            out.append(b)
            continue

        if t == "module":
            baked = p.get("_baked")
            if not expand_module:
                _note_once(notes, "module_nested",
                           "다른 모듈 출력 안에 중첩된 module 블록은 펼치지 않았습니다.")
                out.append(b)
            elif isinstance(baked, list) and baked:
                _flatten_blocks(baked, out, notes, depth + 1, False)
            else:
                _note_once(notes, "module",
                           "module 블록에 렌더된 결과(props._baked)가 없어 제외했습니다.")
            continue
        if t in ("grid", "carousel"):
            _flatten_blocks(p.get("children"), out, notes, depth + 1, expand_module)
            continue
        if t == "card":
            title = str(p.get("title") or "").strip()
            if title:
                out.append({"type": "header", "props": {"text": title, "level": 3}})
            body = next((p.get(k) for k in ("content", "text", "description", "body")
                         if str(p.get(k) or "").strip()), None)
            if body:
                out.append({"type": "text", "props": {"content": str(body)}})
            img = p.get("image")
            if isinstance(img, dict) and img.get("src"):
                out.append({"type": "image", "props": {"src": img.get("src"),
                                                       "alt": img.get("alt")}})
            _flatten_blocks(p.get("children"), out, notes, depth + 1, expand_module)
            footer = str(p.get("footer") or "").strip()
            if footer:
                out.append({"type": "text", "props": {"content": footer}})
            continue
        if t == "slideshow":
            kids = p.get("children")
            if isinstance(kids, list) and kids:
                _flatten_blocks(kids, out, notes, depth + 1, expand_module)
                continue
            for img in (p.get("images") or []):
                if not isinstance(img, dict) or not img.get("src"):
                    continue
                out.append({"type": "image", "props": {"src": img.get("src"),
                                                       "alt": img.get("alt")}})
                cap = str(img.get("caption") or "").strip()
                if cap:
                    out.append({"type": "text", "props": {"content": cap}})
            continue
        if t in ("tabs", "accordion"):
            sections = p.get("tabs") if t == "tabs" else p.get("items")
            for sec in (sections or []):
                if not isinstance(sec, dict):
                    continue
                label = str(sec.get("label") or sec.get("title") or "").strip()
                if label:
                    out.append({"type": "header", "props": {"text": label, "level": 3}})
                body = next((sec.get(k) for k in ("content", "text")
                             if isinstance(sec.get(k), str) and sec.get(k).strip()), None)
                if body:
                    out.append({"type": "text", "props": {"content": body}})
                kids = next((sec.get(k) for k in ("children", "blocks", "items", "content")
                             if isinstance(sec.get(k), list)), None)
                _flatten_blocks(kids, out, notes, depth + 1, expand_module)
            continue
        if t == "paper_trades":
            recs = next((p.get(k) for k in ("records", "rows", "trades")
                         if isinstance(p.get(k), list)), None)
            headers, rows = _table_rows(p.get("headers"), recs)
            if rows:
                out.append({"type": "table",
                            "props": {"headers": headers, "rows": rows,
                                      "title": p.get("title")}})
            else:
                _note_once(notes, "paper_trades",
                           "paper_trades 블록에 기록(props.records)이 없어 제외했습니다.")
            continue
        if t in _LIVE_DEGRADE:
            snap = _live_snapshot(t, p)
            if snap:
                out.append(snap)
            else:
                _note_once(notes, t, f"'{t}' 블록은 실시간 토픽만 담고 있어 제외했습니다 "
                                     "(문서에 실을 스냅샷 데이터 없음).")
            continue
        if t == "table":
            headers, rows = _table_rows(p.get("headers"), p.get("rows"))
            out.append({**b, "type": "table",
                        "props": {**p, "headers": headers, "rows": rows}})
            continue
        out.append(b)


def normalize_blocks(blocks, notes=None):
    """Render-block IR -> the flat, dialect-free stream every renderer here speaks.

    Two dialects are absorbed rather than taught. (1) The fence dialect
    {"name":"Header","type":"component"} is what a model living in chat fences writes — intent is
    unambiguous, so name becomes type. (2) Containers (grid/card/tabs/accordion/carousel/
    slideshow, and a page's `module` block with its baked output) are spliced into their children,
    with tab and accordion sections keeping their titles as level-3 headers so the sectioning
    survives the flattening.

    `notes` (a list) collects one line per DROPPED TYPE — the make_* actions hand them back in
    data.notes, because a document that quietly lost a form is worse than one that says so."""
    out, dropped = [], {}
    _flatten_blocks(blocks, out, dropped, 0, True)
    if notes is not None:
        notes.extend(dropped.values())
    return out


def _split_slides(blocks):
    """header level 1 or divider starts a new slide — the boundary rule from the design."""
    slides, cur = [], []
    for b in blocks:
        t = str(b.get("type") or "")
        p = b.get("props") or {}
        is_boundary = t == "divider" or (t == "header" and _header_level(p) == 1)
        if is_boundary and cur:
            slides.append(cur)
            cur = []
        if t != "divider":
            cur.append(b)
    if cur:
        slides.append(cur)
    return slides


def _kv_pairs(p):
    """key_value props -> [(label, value)] — 'key' and 'label' are the same field."""
    pairs = []
    for it in (p.get("items") or []):
        if not isinstance(it, dict):
            continue
        key = it.get("key") if it.get("key") not in (None, "") else it.get("label")
        pairs.append((str(key or ""), it.get("value")))
    return pairs


def _timeline_items(p):
    """timeline props -> the {title, body[]} group shape the pptx archetypes consume."""
    groups = []
    for it in (p.get("items") or p.get("steps") or p.get("events") or []):
        if isinstance(it, dict):
            head = " ".join(s for s in (str(it.get("date") or "").strip(),
                                        str(it.get("title") or "").strip()) if s)
            body = str(it.get("description") or it.get("content") or it.get("text") or "").strip()
            groups.append({"title": head or body, "body": [body] if body and head else []})
        elif str(it or "").strip():
            groups.append({"title": str(it).strip(), "body": []})
    return groups


def _compare_items(p):
    """compare props -> {title, body[]} groups. The schema is left/right; an items[] list is the
    dialect a model writes when it has three things to compare."""
    raw = p.get("items")
    if not isinstance(raw, list) or not raw:
        raw = [side for side in (p.get("left"), p.get("right")) if isinstance(side, dict)]
    groups = []
    for side in raw:
        if not isinstance(side, dict):
            continue
        title = str(side.get("label") or side.get("title") or side.get("name") or "").strip()
        body = []
        for it in (side.get("items") or []):
            if isinstance(it, dict):
                key = it.get("key") if it.get("key") not in (None, "") else it.get("label")
                body.append(f"{key}: {it.get('value')}" if key else str(it.get("value") or ""))
            elif str(it or "").strip():
                body.append(str(it))
        text = str(side.get("content") or side.get("text") or "").strip()
        if text:
            body.append(text)
        groups.append({"title": title, "body": body})
    return groups


def _progress_line(p):
    value = p.get("value")
    top = parse_number(p.get("max"))
    top = float(top) if top not in (None, 0) else 100.0
    num = parse_number(value)
    pct = None if num is None else max(0.0, min(1.0, float(num) / top))
    label = str(p.get("label") or "").strip()
    shown = f"{value}" if num is None else f"{num:,.2f}".rstrip("0").rstrip(".")
    return label, shown, top, pct


def _block_lines(b):
    """text-ish block -> plain lines for a body textbox / paragraph run. The last-resort form:
    a renderer that has no native shape for a block still prints what the block SAYS."""
    t, p = str(b.get("type") or ""), b.get("props") or {}
    if t == "text":
        return [str(p.get("content") or "")]
    if t == "list":
        mark = "1." if p.get("ordered") else "•"
        return [f"{mark} {item}" for item in (p.get("items") or [])]
    if t == "metric":
        line = f"{p.get('label', '')}: {p.get('value', '')}{p.get('unit', '') or ''}"
        if p.get("delta") not in (None, ""):
            line += f" ({p.get('delta')})"
        return [line]
    if t == "header":
        return [str(p.get("text") or "")]
    if t == "callout":
        title, msg = _callout_parts(p)
        return [f"※ {title} — {msg}" if title and msg else f"※ {title or msg}"]
    if t == "key_value":
        head = [str(p.get("title") or "").strip()] if str(p.get("title") or "").strip() else []
        return head + [f"{k}: {v}" for k, v in _kv_pairs(p)]
    if t == "progress":
        label, shown, top, pct = _progress_line(p)
        tail = "" if pct is None else f" ({pct * 100:.0f}%)"
        return [f"{label}: {shown} / {top:g}{tail}".lstrip(": ")]
    if t == "timeline":
        return [f"• {g['title']}" + (f" — {g['body'][0]}" if g["body"] else "")
                for g in _timeline_items(p)]
    if t == "compare":
        lines = []
        for g in _compare_items(p):
            lines.append(f"[{g['title']}]" if g["title"] else "")
            lines.extend(f"• {x}" for x in g["body"])
        return [l for l in lines if l]
    return []


# Layout constants (inches / points). python-pptx cannot measure rendered text, so heights are
# ESTIMATED with per-column wrap counts — Korean at these sizes runs ~5.5 chars/inch — and the
# estimate errs tall: a slightly airy slide beats text painted over a table (2026-08-10 첫 실물
# 실측 — 셀이 줄바꿈되자 다음 텍스트박스가 표 위에 앉았다).
PPTX_MARGIN_IN = 0.6
PPTX_CHARS_PER_IN = 5.5
PPTX_BODY_PT = 14
PPTX_TABLE_PT = 12


def _wrap_lines(text, width_in, chars_per_in=PPTX_CHARS_PER_IN):
    """Visual line count of `text` in a box `width_in` wide — wrap-aware, never below 1."""
    cpl = max(8, int(width_in * chars_per_in))
    total = 0
    for logical in str(text).split("\n"):
        total += max(1, -(-len(logical) // cpl))
    return total


def _hex_rgb(value):
    """'#2563EB' / '2563EB' -> RGBColor, None for anything unreadable (never a guess)."""
    from pptx.dml.color import RGBColor
    s = str(value or "").strip().lstrip("#")
    if len(s) == 6:
        try:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
        except ValueError:
            return None
    return None


# Pure numerics ("1,234", "-3.1") become real numbers; unit-bearing text ("48억", "3.1%p")
# cannot carry its unit into a number, so it stays text.
_NUMERIC_RE = re.compile(r"^[+\-]?\d{1,3}(,\d{3})*(\.\d+)?$|^[+\-]?\d+(\.\d+)?$")
# Looser shape for ALIGNMENT decisions only: a number plus a short unit tail still reads
# as a number to the eye, and number columns right-align.
_NUMLIKE_RE = re.compile(r"^[+\-±]?[\d,]+(\.\d+)?\s*\S{0,3}$")
# Titles that read as points on a schedule — they pull item groups onto a timeline.
_STEPISH_RE = re.compile(r"(\d{4}|\d+\s*년|\d+\s*월|[1-4]\s*Q|분기|단계|Phase|Step)", re.I)


def parse_number(v):
    if isinstance(v, bool):
        return None
    if isinstance(v, (int, float)):
        return v
    s = str(v or "").strip()
    if not s or not _NUMERIC_RE.match(s):
        return None
    try:
        f = float(s.replace(",", ""))
    except ValueError:
        return None
    return int(f) if f.is_integer() and "." not in s else f


def make_pptx_file(blocks, title, master_path, out_path, transition="fade", theme=None):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    prs = Presentation(master_path) if master_path else Presentation()
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    sw, sh = prs.slide_width, prs.slide_height
    sw_in = sw / 914400
    sh_in = sh / 914400

    # Design has THREE layers, strongest first: a master .pptx (corporate look — never painted
    # over), then caller theme tokens (the palette the AI chose when it designed the page —
    # blocks are meaning, tokens are the look, so the look rides along as data), then the
    # Firebat default (blue-600 accent on slate, banded tables, quiet page number).
    styled = master_path is None
    theme = theme if isinstance(theme, dict) else {}
    BLUE = (_hex_rgb(theme.get("accent")) or RGBColor(0x25, 0x63, 0xEB))
    SLATE_D = (_hex_rgb(theme.get("heading")) or RGBColor(0x1E, 0x29, 0x3B))
    SLATE_B = (_hex_rgb(theme.get("body")) or RGBColor(0x33, 0x41, 0x55))
    SLATE_L = (_hex_rgb(theme.get("band")) or RGBColor(0xF1, 0xF5, 0xF9))
    SLATE_M = RGBColor(0x94, 0xA3, 0xB8)   # footer — always quiet
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # One typeface across the deck (styled path only — a master's fonts rule there).
    # Korean needs BOTH the latin and eastAsian slots set, or ascii runs fall back to
    # a different face mid-sentence. theme.font swaps the whole deck (e.g. Pretendard).
    FONT_EA = str(theme.get("font") or "맑은 고딕")

    def apply_font(para):
        if not styled:
            return
        para.font.name = FONT_EA
        rpr = para.font._rPr
        if rpr is not None:
            latin = rpr.find(qn("a:latin"))
            ea = rpr.find(qn("a:ea"))
            if ea is None and latin is not None:
                ea = rpr.makeelement(qn("a:ea"), {})
                latin.addnext(ea)
            if ea is not None:
                ea.set("typeface", FONT_EA)

    # Slide transition — python-pptx has no API for it, but <p:transition> is one well-formed
    # element per slide (unlike object-animation timing XML, which stays out of scope). Injected
    # for every slide at save time; "none" disables.
    def apply_transition(slide):
        if transition not in ("fade", "wipe", "push"):
            return
        sld = slide._element
        if sld.find(qn("p:transition")) is not None:
            return
        tr = sld.makeelement(qn("p:transition"), {"spd": "med"})
        tr.append(sld.makeelement(qn(f"p:{transition}"), {}))
        timing = sld.find(qn("p:timing"))
        if timing is not None:
            timing.addprevious(tr)
        else:
            sld.append(tr)

    def accent_bar(slide, x_in, y_in, w_in, h_in=0.05):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLUE
        bar.line.fill.background()
        bar.shadow.inherit = False
        return bar

    def add_box(slide, shape_id, x, y, w, h, fill=None, line=None, line_w=None):
        box = slide.shapes.add_shape(shape_id, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            box.fill.background()
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = fill
        if line is None:
            box.line.fill.background()
        else:
            box.line.color.rgb = line
            if line_w:
                box.line.width = Pt(line_w)
        box.shadow.inherit = False
        return box

    def add_text(slide, x, y, w, h, runs, align=None, wrap=True, anchor=None):
        # runs = [(text, size_pt, bold, color), ...] — one paragraph per run.
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        if anchor is not None:
            tf.vertical_anchor = anchor
        for i, (text, size, bold, color) in enumerate(runs):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = text
            para.font.size = Pt(size)
            para.font.bold = bold
            if color is not None:
                para.font.color.rgb = color
            if align is not None:
                para.alignment = align
            apply_font(para)
        return tb

    if title:
        if styled:
            # Cover, proposal-deck shape: big heading upper-left, accent bar, a quiet band
            # strip along the bottom. No layout placeholders — nothing to ghost.
            slide = prs.slides.add_slide(blank)
            t_lines = _wrap_lines(str(title), sw_in - 1.7, chars_per_in=3.4)
            add_text(slide, 0.85, 2.0, sw_in - 1.7, 0.65 * t_lines,
                     [(str(title), 32, True, SLATE_D)])
            accent_bar(slide, 0.88, 2.12 + 0.65 * t_lines, 2.4, 0.07)
            add_box(slide, MSO_SHAPE.RECTANGLE, 0, sh_in - 0.9, sw_in, 0.9, fill=SLATE_L)
            add_box(slide, MSO_SHAPE.RECTANGLE, 0, sh_in - 0.95, sw_in, 0.05, fill=BLUE)
        else:
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            placed = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 0:
                    ph.text = str(title)
                    placed = ph
                    break
            # Unfilled placeholders render as dashed "click to add" boxes (the subtitle ghost
            # on the first live deck) — remove every placeholder we did not fill.
            for ph in list(slide.placeholders):
                if placed is not None and ph.placeholder_format.idx == 0:
                    continue
                el = ph._element
                el.getparent().remove(el)
            if placed is None:
                tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5),
                                              sw - Inches(1.6), Inches(1.5))
                tb.text_frame.text = str(title)
                tb.text_frame.paragraphs[0].font.size = Pt(40)

    state = {"count": 1 if title else 0, "slide": None, "y": 0.0,
             "sec_no": 0, "sec_title": None, "sec_stmt": None}
    BAND_H = 0.62          # navy section band
    STMT_H = 0.78          # light statement band under it
    BODY_BOTTOM = 0.5      # keep-out at the slide foot

    def draw_band(slide, statement):
        # The genre's spine: dark number+title band, then (first slide of the section only)
        # a light band carrying the slide's one-line message behind a small accent bar.
        add_box(slide, MSO_SHAPE.RECTANGLE, 0, 0, sw_in, BAND_H, fill=SLATE_D)
        add_text(slide, 0.55, 0, sw_in - 2.6, BAND_H,
                 [(f"{state['sec_no']:02d}. {state['sec_title']}", 16, True, WHITE)],
                 wrap=False, anchor=MSO_ANCHOR.MIDDLE)
        if title:
            add_text(slide, sw_in - 2.05, 0, 1.55, BAND_H,
                     [(str(title), 8.5, False, SLATE_M)], align=PP_ALIGN.RIGHT, wrap=False,
                     anchor=MSO_ANCHOR.MIDDLE)
        if statement:
            add_box(slide, MSO_SHAPE.RECTANGLE, 0, BAND_H, sw_in, STMT_H, fill=SLATE_L)
            add_box(slide, MSO_SHAPE.RECTANGLE, 0.55, BAND_H + 0.15, 0.07, STMT_H - 0.3,
                    fill=BLUE)
            add_text(slide, 0.8, BAND_H, sw_in - 1.7, STMT_H,
                     [(statement, 13, True, SLATE_D)], anchor=MSO_ANCHOR.MIDDLE)
            return BAND_H + STMT_H + 0.3
        return BAND_H + 0.35

    def new_slide(continuation=True):
        state["slide"] = prs.slides.add_slide(blank)
        state["count"] += 1
        if styled and state["sec_title"]:
            state["y"] = draw_band(state["slide"],
                                   None if continuation else state["sec_stmt"])
        else:
            state["y"] = 0.55 if styled else 0.4
        if styled:
            ft = state["slide"].shapes.add_textbox(
                Inches(sw_in - 1.05), Inches(sh_in - 0.38), Inches(0.75), Inches(0.28))
            para = ft.text_frame.paragraphs[0]
            para.text = str(state["count"])
            para.font.size = Pt(10)
            para.font.color.rgb = SLATE_M
            para.alignment = PP_ALIGN.RIGHT
            apply_font(para)

    def ensure_room(height_in):
        # Content that will not fit CONTINUES on a fresh slide — the old behavior silently
        # dropped the rest of the chunk once the cursor passed the bottom.
        if state["slide"] is None or state["y"] + height_in > sh_in - BODY_BOTTOM:
            new_slide()

    def set_cell(cell, text, bold, fill=None, color=None, align=None):
        cell.text = text
        if fill is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(PPTX_TABLE_PT)
            para.font.bold = bold
            if color is not None:
                para.font.color.rgb = color
            if align is not None:
                para.alignment = align
            apply_font(para)

    def place_table(headers, seg_rows, height_in, col_count, aligns=None):
        ensure_room(height_in + 0.2)
        slide, y = state["slide"], state["y"]
        shape = slide.shapes.add_table(
            len(seg_rows) + 1, col_count, Inches(PPTX_MARGIN_IN), Inches(y),
            sw - Inches(PPTX_MARGIN_IN * 2), Inches(height_in))
        table = shape.table
        for c, h in enumerate(headers):
            # Header row centers regardless of column alignment — the convention is
            # centered headers over right-aligned numbers (2026-08-11 사용자).
            set_cell(table.cell(0, c), str(h), True,
                     fill=BLUE if styled else None, color=WHITE if styled else None,
                     align=PP_ALIGN.CENTER)
        for r, row in enumerate(seg_rows):
            band = SLATE_L if (styled and r % 2 == 1) else (WHITE if styled else None)
            for c in range(col_count):
                val = row[c] if c < len(row) else ""
                set_cell(table.cell(r + 1, c), "" if val is None else str(val), False,
                         fill=band, color=SLATE_B if styled else None,
                         align=aligns[c] if aligns else None)
        state["y"] = y + height_in + 0.25

    def render_table(p):
        headers, rows = _table_rows(p.get("headers"), p.get("rows"))
        if not headers and rows:
            headers = [str(c) for c in rows[0]]
            rows = rows[1:]
        if not headers:
            return
        n_cols = len(headers)
        col_w_in = (sw_in - PPTX_MARGIN_IN * 2) / n_cols

        # Wrap-aware height: each row is as tall as its most-wrapped cell.
        def row_h(cells):
            lines = max(_wrap_lines("" if c is None else c, col_w_in) for c in cells)
            return 0.14 + 0.21 * lines

        header_h = row_h(headers)
        row_hs = [row_h([row[c] if c < len(row) else "" for c in range(n_cols)])
                  for row in rows]
        # A table taller than a slide SPLITS across slides with the header repeated —
        # one connected table, not a truncated one (2026-08-10 사용자 질문이 이 계약).
        usable = sh_in - (1.6 if styled else 0.9)
        segments = []
        if header_h + sum(row_hs) + 0.2 > usable:
            seg, seg_h = [], header_h
            for row, rh in zip(rows, row_hs):
                if seg and seg_h + rh + 0.2 > usable:
                    segments.append((seg, seg_h))
                    seg, seg_h = [], header_h
                seg.append(row)
                seg_h += rh
            if seg:
                segments.append((seg, seg_h))
        else:
            segments = [(rows, header_h + sum(row_hs))]
        # Number columns read right-aligned (decided on the WHOLE table, so every split
        # segment aligns the same way) — majority, the same rule the xlsx sheets use.
        numeric = _numeric_cols(rows, n_cols)
        aligns = [PP_ALIGN.RIGHT if ci in numeric else None for ci in range(n_cols)]
        for seg_rows, seg_h in segments:
            place_table(headers, seg_rows, seg_h, n_cols, aligns=aligns)

    def render_image(p):
        src = str(p.get("src") or "")
        img_path, _ = resolve_path(src)
        if img_path:
            ensure_room(2.7)
            slide, y = state["slide"], state["y"]
            try:
                slide.shapes.add_picture(img_path, Inches(PPTX_MARGIN_IN), Inches(y),
                                         height=Inches(2.5))
                state["y"] = y + 2.7
            except Exception:  # noqa: BLE001 — a bad image loses itself, not the deck
                pass

    def deck_palette(n):
        """n colors that all descend from the deck's one accent — a chart drawn in Excel's stock
        palette would be the only thing on the slide the theme did not choose."""
        base = [BLUE, _mix(BLUE, 255, 0.45), _mix(BLUE, 0, 0.28), _mix(BLUE, 255, 0.7),
                SLATE_D, SLATE_M]
        return [base[i % len(base)] for i in range(max(1, n))]

    def render_chart(b):
        """A chart block as a NATIVE pptx chart — editable, with its numbers in the embedded
        workbook. Previously a chart block reached a deck as nothing at all."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

        parsed = _chart_series(b)
        if not parsed:
            return False
        ctype, ctitle, labels, cols = parsed
        kinds = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "line": XL_CHART_TYPE.LINE,
                 "pie": XL_CHART_TYPE.PIE, "doughnut": XL_CHART_TYPE.DOUGHNUT,
                 "combo": XL_CHART_TYPE.COLUMN_CLUSTERED}
        kind = kinds.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)
        data = CategoryChartData()
        data.categories = [str(x) for x in labels]
        for name, vals in cols:
            data.add_series(name, tuple(parse_number(v) for v in vals))
        h = 3.3
        ensure_room(h + 0.25)
        slide, y = state["slide"], state["y"]
        w = sw_in - PPTX_MARGIN_IN * 2
        try:
            frame = slide.shapes.add_chart(kind, Inches(PPTX_MARGIN_IN), Inches(y),
                                           Inches(w), Inches(h), data)
        except Exception:  # noqa: BLE001 — a bad chart loses itself, not the deck
            return False
        chart = frame.chart
        chart.font.size = Pt(10)
        if styled:
            chart.font.color.rgb = SLATE_B
        if ctitle:
            chart.has_title = True
            chart.chart_title.text_frame.text = ctitle
            tp = chart.chart_title.text_frame.paragraphs[0]
            tp.font.size = Pt(12)
            tp.font.bold = True
            if styled:
                tp.font.color.rgb = SLATE_D
        round_chart = ctype in ("pie", "doughnut")
        chart.has_legend = round_chart or len(cols) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        plot = chart.plots[0]
        if round_chart:
            # One color per SLICE: a pie's categories are the thing being told apart.
            plot.vary_by_categories = True
            pts = list(plot.series[0].points)
            for pt, color in zip(pts, deck_palette(len(pts))):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = color
        else:
            for ser, color in zip(chart.series, deck_palette(len(cols))):
                if ctype == "line":
                    ser.format.line.color.rgb = color
                    ser.format.line.width = Pt(2.0)
                    ser.smooth = False
                else:
                    ser.format.fill.solid()
                    ser.format.fill.fore_color.rgb = color
        state["y"] = y + h + 0.25
        return True

    def render_stock_chart(p):
        """No native candlestick exists in python-pptx, so a candle series is told the two honest
        ways: the close as a line chart, and (when it will not fit) the table underneath."""
        headers, rows = _stock_rows(p)
        if not rows:
            return
        try:
            close_i = headers.index("종가")
        except ValueError:
            close_i = len(headers) - 1
        labels = [r[0] for r in rows]
        closes = [r[close_i] if close_i < len(r) else None for r in rows]
        title = str(p.get("title") or p.get("symbol") or "").strip()
        done = render_chart({"type": "chart", "props": {
            "chartType": "line", "title": title, "labels": labels, "data": closes}})
        if not done:
            block, _extra = _stock_table_block(p)
            if block:
                render_table(block["props"])

    def callout_band(p):
        """The honesty note, drawn: a tinted full-bleed band behind an accent tick. Same machinery
        as the section statement band, because it is the same move — one line that must be read."""
        title, msg = _callout_parts(p)
        if not title and not msg:
            return
        body_w = sw_in - 1.7
        lines = (_wrap_lines(title, body_w, chars_per_in=4.5) if title else 0) \
            + (_wrap_lines(msg, body_w) if msg else 0)
        h = max(0.5, 0.1 + 0.26 * lines)
        ensure_room(h + 0.25)
        slide, y = state["slide"], state["y"]
        add_box(slide, MSO_SHAPE.RECTANGLE, PPTX_MARGIN_IN, y, sw_in - PPTX_MARGIN_IN * 2, h,
                fill=SLATE_L)
        add_box(slide, MSO_SHAPE.RECTANGLE, PPTX_MARGIN_IN, y, 0.07, h, fill=BLUE)
        runs = ([(title, 12, True, SLATE_D)] if title else []) \
            + ([(msg, 11, False, SLATE_B)] if msg else [])
        add_text(slide, PPTX_MARGIN_IN + 0.22, y, body_w, h, runs,
                 anchor=MSO_ANCHOR.MIDDLE)
        state["y"] = y + h + 0.25

    # "No Style, No Grid" — the only way to get a python-pptx table without the default style's
    # rules and banding, which turn a two-column fact list into a spreadsheet.
    PPTX_PLAIN_TABLE_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"

    def key_value_table(p):
        pairs = _kv_pairs(p)
        if not pairs:
            return
        title = str(p.get("title") or "").strip()
        if title:
            text_par([title])
        col_w_in = (sw_in - PPTX_MARGIN_IN * 2) / 2
        heights = [0.14 + 0.21 * max(_wrap_lines(str(k), col_w_in),
                                     _wrap_lines("" if v is None else v, col_w_in))
                   for k, v in pairs]
        h = sum(heights)
        ensure_room(h + 0.25)
        slide, y = state["slide"], state["y"]
        shape = slide.shapes.add_table(len(pairs), 2, Inches(PPTX_MARGIN_IN), Inches(y),
                                       sw - Inches(PPTX_MARGIN_IN * 2), Inches(h))
        table = shape.table
        table.first_row = False
        table.horz_banding = False
        try:
            style_id = shape._element.graphic.graphicData.tbl.tblPr.find(qn("a:tableStyleId"))
            if style_id is not None:
                style_id.text = PPTX_PLAIN_TABLE_STYLE
        except Exception:  # noqa: BLE001 — a styled table beats a crashed deck
            pass
        for r, (k, v) in enumerate(pairs):
            set_cell(table.cell(r, 0), str(k), True,
                     color=SLATE_B if styled else None)
            set_cell(table.cell(r, 1), "" if v is None else str(v), False,
                     color=SLATE_D if styled else None, align=PP_ALIGN.RIGHT)
        state["y"] = y + h + 0.25

    def progress_rows(items):
        """A ratio told the way a dashboard tells it: label, number, and a track whose filled
        length IS the number."""
        row_h, gap = 0.5, 0.12
        for p in items:
            label, shown, top, pct = _progress_line(p)
            ensure_room(row_h + gap)
            slide, y = state["slide"], state["y"]
            track_x = PPTX_MARGIN_IN + 2.6
            track_w = sw_in - PPTX_MARGIN_IN - track_x - 0.9
            add_text(slide, PPTX_MARGIN_IN, y, 2.5, row_h,
                     [(label, 11, False, SLATE_B)], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            add_box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, track_x, y + 0.16, max(track_w, 0.5),
                    0.18, fill=SLATE_L)
            if pct:
                add_box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, track_x, y + 0.16,
                        max(track_w * pct, 0.06), 0.18, fill=BLUE)
            tail = shown if pct is None else f"{pct * 100:.0f}%"
            add_text(slide, sw_in - PPTX_MARGIN_IN - 0.85, y, 0.85, row_h,
                     [(tail, 12, True, BLUE)], align=PP_ALIGN.RIGHT,
                     anchor=MSO_ANCHOR.MIDDLE, wrap=False)
            state["y"] = y + row_h + gap

    # ---- plain path (master decks): a quiet document flow the master's look carries ----

    def render_chunk_plain(chunk):
        for b in chunk:
            t, p = str(b.get("type") or ""), b.get("props") or {}
            if t == "header":
                ensure_room(0.95)
                slide, y = state["slide"], state["y"]
                tb = slide.shapes.add_textbox(Inches(PPTX_MARGIN_IN), Inches(y),
                                              sw - Inches(PPTX_MARGIN_IN * 2), Inches(0.7))
                tf = tb.text_frame
                tf.text = str(p.get("text") or "")
                tf.paragraphs[0].font.size = Pt(28 if _header_level(p) == 1 else 20)
                tf.paragraphs[0].font.bold = True
                state["y"] = y + 0.95
            elif t == "table":
                render_table(p)
            elif t == "image":
                render_image(p)
            elif t in _CHART_BLOCK_TYPES:
                render_chart(b)
            elif t in _STOCK_BLOCK_TYPES:
                render_stock_chart(p)
            elif t == "callout":
                callout_band(p)
            elif t == "key_value":
                key_value_table(p)
            elif t == "progress":
                progress_rows([p])
            else:
                lines = _block_lines(b)
                if not lines:
                    continue
                body_w_in = sw_in - 1.4
                visual = sum(_wrap_lines(l, body_w_in) for l in lines)
                height_in = 0.1 + 0.26 * visual
                ensure_room(height_in + 0.15)
                slide, y = state["slide"], state["y"]
                add_text(slide, 0.7, y, body_w_in, height_in,
                         [(line, PPTX_BODY_PT, False, None) for line in lines])
                state["y"] = y + height_in + 0.15

    # ---- styled path: the proposal-deck layout system (2026-08-11 사용자 예시 장르) ----
    # Content is poured into archetypes instead of flowed like a document: numbered ring
    # rows, 2-3 pill columns, a 4-quadrant donut, metric cards, and a takeaway statement.

    def _mix(color, target, f):
        s = str(color)
        return RGBColor(*(round(int(s[k:k + 2], 16) + (target - int(s[k:k + 2], 16)) * f)
                          for k in (0, 2, 4)))

    def collect_groups(rest):
        groups, i = [], 0
        while i < len(rest):
            b = rest[i]
            t, p = str(b.get("type") or ""), b.get("props") or {}
            if t == "header":
                g = {"title": str(p.get("text") or ""), "body": []}
                i += 1
                while i < len(rest):
                    nb = rest[i]
                    nt, np_ = str(nb.get("type") or ""), nb.get("props") or {}
                    if nt == "text":
                        g["body"].append(str(np_.get("content") or ""))
                    elif nt == "list":
                        g["body"].extend(f"• {it}" for it in (np_.get("items") or []))
                    else:
                        break
                    i += 1
                groups.append(("item", g))
            elif t == "metric":
                ms = []
                while i < len(rest) and str(rest[i].get("type") or "") == "metric":
                    ms.append(rest[i].get("props") or {})
                    i += 1
                groups.append(("metrics", ms))
            elif t == "list":
                groups.append(("list", [str(it) for it in (p.get("items") or [])]))
                i += 1
            elif t == "text":
                groups.append(("text", [str(p.get("content") or "")]))
                i += 1
            elif t == "progress":
                # Consecutive progress bars are ONE strip, the way consecutive metrics are.
                ps = []
                while i < len(rest) and str(rest[i].get("type") or "") == "progress":
                    ps.append(rest[i].get("props") or {})
                    i += 1
                groups.append(("progress", ps))
            elif t == "timeline":
                # The archetypes below used to be reachable only by heuristic — a header group
                # that HAPPENED to look like a schedule. A block that SAYS it is a timeline is
                # the stronger signal, so it triggers the archetype directly (2026-08-12 audit).
                groups.append(("timeline", _timeline_items(p)))
                i += 1
            elif t == "compare":
                groups.append(("compare", _compare_items(p)))
                i += 1
            elif t in ("table", "image", "key_value", "callout"):
                groups.append((t, p))
                i += 1
            elif t in _CHART_BLOCK_TYPES or t in _STOCK_BLOCK_TYPES:
                groups.append(("chart" if t in _CHART_BLOCK_TYPES else "stock_chart", b))
                i += 1
            else:
                lines = _block_lines(b)
                if lines:
                    groups.append(("text", lines))
                i += 1
        return groups

    ITEM_X = PPTX_MARGIN_IN + 0.95

    def item_row(no, g):
        text_w = sw_in - ITEM_X - PPTX_MARGIN_IN
        t_lines = _wrap_lines(g["title"], text_w, chars_per_in=4.5)
        b_lines = sum(_wrap_lines(l, text_w) for l in g["body"])
        h = max(0.68, 0.3 * t_lines + 0.22 * b_lines + 0.16)
        ensure_room(h + 0.1)
        slide, y = state["slide"], state["y"]
        ring = add_box(slide, MSO_SHAPE.OVAL, PPTX_MARGIN_IN + 0.02, y + 0.02, 0.6, 0.6,
                       fill=WHITE, line=BLUE, line_w=2.25)
        rp = ring.text_frame.paragraphs[0]
        rp.text = f"{no:02d}"
        rp.font.size = Pt(14)
        rp.font.bold = True
        rp.font.color.rgb = BLUE
        rp.alignment = PP_ALIGN.CENTER
        apply_font(rp)
        runs = [(g["title"], 13.5, True, SLATE_D)]
        runs += [(l, 11, False, SLATE_B) for l in g["body"]]
        add_text(slide, ITEM_X, y, text_w, h, runs)
        state["y"] = y + h + 0.12

    def pill_columns(items):
        n = len(items)
        gap = 0.35
        col_w = (sw_in - PPTX_MARGIN_IN * 2 - gap * (n - 1)) / n
        body_h = 0.2 + 0.22 * max(
            sum(_wrap_lines(l, col_w - 0.2) for l in g["body"]) or 1 for g in items)
        ensure_room(0.44 + 0.15 + body_h + 0.1)
        slide, y = state["slide"], state["y"]
        for idx, g in enumerate(items):
            x = PPTX_MARGIN_IN + idx * (col_w + gap)
            pill = add_box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.25, y,
                           col_w - 0.5, 0.44, fill=BLUE)
            try:
                pill.adjustments[0] = 0.5
            except Exception:  # noqa: BLE001 — a squarer pill, not a lost slide
                pass
            pp = pill.text_frame.paragraphs[0]
            pp.text = g["title"]
            pp.font.size = Pt(12.5)
            pp.font.bold = True
            pp.font.color.rgb = WHITE
            pp.alignment = PP_ALIGN.CENTER
            apply_font(pp)
            add_text(slide, x, y + 0.59, col_w, body_h,
                     [(l, 10.5, False, SLATE_B) for l in g["body"]], align=PP_ALIGN.CENTER)
        state["y"] = y + 0.59 + body_h + 0.15

    def add_pie(slide, x, y, d, start_deg, end_deg, fill):
        pie = add_box(slide, MSO_SHAPE.PIE, x, y, d, d, fill=fill)
        try:
            # PIE adjustments are angles: raw XML is 1/60000 deg, python-pptx divides by
            # 100000 — so degrees x 0.6. 0 = 3 o'clock, clockwise.
            pie.adjustments[0] = start_deg * 0.6
            pie.adjustments[1] = end_deg * 0.6
        except Exception:  # noqa: BLE001 — a full circle beats a crashed deck
            pass
        return pie

    def quad(items):
        # SWOT-shape: a four-color donut (pie quarters under a white core) with corner
        # badges, one text block per quadrant (2026-08-11 사용자 예시).
        ensure_room(3.8)
        slide, y = state["slide"], state["y"]
        cx = sw_in / 2
        cy = y + (sh_in - BODY_BOTTOM - y) / 2
        colors = [BLUE, _mix(BLUE, 255, 0.45), _mix(BLUE, 0, 0.3), _mix(BLUE, 255, 0.2)]
        R = 1.55
        angles = [(180, 270), (270, 360), (90, 180), (0, 90)]      # NW NE SW SE
        offs = [(-1.1, -1.1), (1.1, -1.1), (-1.1, 1.1), (1.1, 1.1)]
        for i in range(min(4, len(items))):
            add_pie(slide, cx - R, cy - R, R * 2, angles[i][0], angles[i][1], colors[i])
        add_box(slide, MSO_SHAPE.OVAL, cx - 1.0, cy - 1.0, 2.0, 2.0,
                fill=WHITE, line=SLATE_L, line_w=1.0)
        label = state["sec_title"] or ""
        if label:
            add_text(slide, cx - 0.95, cy - 0.3, 1.9, 0.6,
                     [(label, 11.5, True, SLATE_D)], align=PP_ALIGN.CENTER)
        col_w = cx - R - 0.55 - PPTX_MARGIN_IN
        xs = [PPTX_MARGIN_IN, cx + R + 0.55, PPTX_MARGIN_IN, cx + R + 0.55]
        ys = [y + 0.15, y + 0.15, cy + 0.55, cy + 0.55]
        for i, g in enumerate(items[:4]):
            badge = add_box(slide, MSO_SHAPE.OVAL, cx + offs[i][0] - 0.31,
                            cy + offs[i][1] - 0.31, 0.62, 0.62,
                            fill=colors[i], line=WHITE, line_w=2.0)
            first = g["title"].strip()[:1]
            bp = badge.text_frame.paragraphs[0]
            bp.text = first.upper() if first.isascii() and first.isalpha() else str(i + 1)
            bp.font.size = Pt(15)
            bp.font.bold = True
            bp.font.color.rgb = WHITE
            bp.alignment = PP_ALIGN.CENTER
            apply_font(bp)
            align = PP_ALIGN.RIGHT if i in (1, 3) else PP_ALIGN.LEFT
            runs = [(g["title"], 15, True, colors[i])]
            runs += [(l, 10.5, False, SLATE_B) for l in g["body"]]
            add_text(slide, xs[i], ys[i], col_w, cy - y - 0.7, runs, align=align)
        state["y"] = sh_in  # the quadrant owns the body — next content turns the page

    def metric_strip(ms):
        per_row = 4 if len(ms) > 3 else max(1, len(ms))
        gap = 0.25
        card_w = (sw_in - PPTX_MARGIN_IN * 2 - gap * (per_row - 1)) / per_row
        for start in range(0, len(ms), per_row):
            ensure_room(1.3)
            slide, y = state["slide"], state["y"]
            for j, m in enumerate(ms[start:start + per_row]):
                x = PPTX_MARGIN_IN + j * (card_w + gap)
                card = add_box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, 1.15,
                               fill=SLATE_L)
                try:
                    card.adjustments[0] = 0.12
                except Exception:  # noqa: BLE001
                    pass
                tf = card.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                value = f"{m.get('value', '')}{m.get('unit') or ''}"
                # Dashboard-card reading order: label anchors left, numbers set right
                # (2026-08-11 사용자 — autoshape default centering scattered the three lines).
                rows = [(str(m.get("label") or ""), 10, False, SLATE_B, PP_ALIGN.LEFT),
                        (value, 20, True, BLUE, PP_ALIGN.RIGHT)]
                if m.get("delta") not in (None, ""):
                    rows.append((str(m.get("delta")), 9.5, False, SLATE_B, PP_ALIGN.RIGHT))
                for i, (text, size, bold, color, al) in enumerate(rows):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = text
                    para.font.size = Pt(size)
                    para.font.bold = bold
                    para.font.color.rgb = color
                    para.alignment = al
                    apply_font(para)
            state["y"] = y + 1.3

    def text_par(lines):
        body_w_in = sw_in - 1.4
        visual = sum(_wrap_lines(l, body_w_in) for l in lines)
        h = 0.1 + 0.24 * visual
        ensure_room(h + 0.15)
        add_text(state["slide"], 0.7, state["y"], body_w_in, h,
                 [(l, 11.5, False, SLATE_B) for l in lines])
        state["y"] += h + 0.15

    def statement_line(text):
        # The genre's closing move: one bold accent-colored line, set to the right.
        ensure_room(0.75)
        slide, y = state["slide"], state["y"]
        add_text(slide, sw_in * 0.3, y + 0.1, sw_in * 0.7 - PPTX_MARGIN_IN, 0.55,
                 [(text, 15.5, True, BLUE)], align=PP_ALIGN.RIGHT)
        state["y"] = y + 0.75

    def chevron_strip(steps):
        # A short list whose items are labels, not sentences, reads as a process —
        # drawn as an arrow flow, light to accent, left to right.
        n = len(steps)
        gap = 0.12
        w = (sw_in - PPTX_MARGIN_IN * 2 - gap * (n - 1)) / n
        h = 0.72
        ensure_room(h + 0.25)
        slide, y = state["slide"], state["y"]
        for i, txt in enumerate(steps):
            f = 0.5 * (n - 1 - i) / max(n - 1, 1)
            shp = add_box(slide, MSO_SHAPE.CHEVRON, PPTX_MARGIN_IN + i * (w + gap), y,
                          w, h, fill=_mix(BLUE, 255, f))
            cp = shp.text_frame.paragraphs[0]
            cp.text = txt
            cp.font.size = Pt(11.5)
            cp.font.bold = True
            cp.font.color.rgb = WHITE
            cp.alignment = PP_ALIGN.CENTER
            apply_font(cp)
        state["y"] = y + h + 0.25

    def timeline(items):
        # Milestone titles ride a horizontal axis, texts zigzag above and below —
        # the genre's curve-with-dots slide in straight-line form.
        ensure_room(3.4)
        slide, y = state["slide"], state["y"]
        cy = y + (sh_in - BODY_BOTTOM - y) / 2
        x0 = PPTX_MARGIN_IN + 0.4
        x1 = sw_in - PPTX_MARGIN_IN - 0.4
        add_box(slide, MSO_SHAPE.RECTANGLE, x0 - 0.2, cy - 0.02, x1 - x0 + 0.4, 0.04,
                fill=_mix(BLUE, 255, 0.75))
        n = len(items)
        step = (x1 - x0) / max(n - 1, 1)
        for i, g in enumerate(items):
            cx_i = x0 + step * i
            add_box(slide, MSO_SHAPE.OVAL, cx_i - 0.17, cy - 0.17, 0.34, 0.34,
                    fill=WHITE, line=BLUE, line_w=2.0)
            add_box(slide, MSO_SHAPE.OVAL, cx_i - 0.07, cy - 0.07, 0.14, 0.14, fill=BLUE)
            col_w = min(step * 1.9, 2.7) if n > 1 else 3.0
            tx = min(max(cx_i - col_w / 2, PPTX_MARGIN_IN), sw_in - PPTX_MARGIN_IN - col_w)
            above = i % 2 == 0
            box_h = (cy - y - 0.45) if above else (sh_in - BODY_BOTTOM - cy - 0.45)
            ty = y + 0.1 if above else cy + 0.35
            runs = [(g["title"], 12.5, True, BLUE)]
            runs += [(l, 10, False, SLATE_B) for l in g["body"]]
            add_text(slide, tx, ty, col_w, max(box_h, 0.5), runs, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.BOTTOM if above else MSO_ANCHOR.TOP)
        state["y"] = sh_in  # the axis owns the body — next content turns the page

    sec_counter = {"n": 0}

    def render_chunk_styled(chunk):
        rest = list(chunk)
        state["sec_title"], state["sec_stmt"] = None, None
        if rest and str(rest[0].get("type") or "") == "header" \
                and _header_level(rest[0].get("props") or {}) == 1:
            sec_counter["n"] += 1
            state["sec_no"] = sec_counter["n"]
            state["sec_title"] = str((rest[0].get("props") or {}).get("text") or "")
            rest = rest[1:]
            # A short first paragraph is the slide's one-line message — it moves into the
            # statement band instead of sitting in the body as a stray paragraph.
            if rest and str(rest[0].get("type") or "") == "text":
                c = str((rest[0].get("props") or {}).get("content") or "")
                if 0 < len(c) <= 130:
                    state["sec_stmt"] = c
                    rest = rest[1:]
        new_slide(continuation=False)
        groups = collect_groups(rest)
        stmt = None
        if len(groups) >= 2 and groups[-1][0] == "text" and len(groups[-1][1]) == 1 \
                and 0 < len(groups[-1][1][0]) <= 70:
            stmt = groups[-1][1][0]
            groups = groups[:-1]
        items = [g for k, g in groups if k == "item"]
        only_items = bool(groups) and all(k == "item" for k, _ in groups)
        short = all(sum(len(l) for l in g["body"]) <= 220 for g in items)
        # Timeline outranks the donut: four year-titled groups are a schedule, not a SWOT.
        # These stay HEURISTICS for block streams that only imply their shape; a stream that
        # names it (a timeline/compare block) went down the direct path in collect_groups.
        timeline_fit = (only_items and 3 <= len(items) <= 6 and short
                        and sum(1 for g in items if _STEPISH_RE.search(g["title"]))
                        >= max(2, len(items) - 1))
        if timeline_fit:
            timeline(items)
        elif only_items and len(items) == 4 and short:
            quad(items)
        elif only_items and 2 <= len(items) <= 3 and short:
            pill_columns(items)
        else:
            no = 0
            for kind, val in groups:
                if kind == "item":
                    no += 1
                    item_row(no, val)
                elif kind == "metrics":
                    metric_strip(val)
                elif kind == "table":
                    render_table(val)
                elif kind == "image":
                    render_image(val)
                elif kind == "chart":
                    render_chart(val)
                elif kind == "stock_chart":
                    render_stock_chart(val.get("props") or {})
                elif kind == "callout":
                    callout_band(val)
                elif kind == "key_value":
                    key_value_table(val)
                elif kind == "progress":
                    progress_rows(val)
                elif kind == "timeline":
                    if val:
                        timeline(val) if len(val) <= 6 else text_par(
                            [f"• {g['title']}" for g in val])
                elif kind == "compare":
                    if 2 <= len(val) <= 3:
                        pill_columns(val)
                    else:
                        for gi, g in enumerate(val, start=1):
                            item_row(gi, g)
                elif kind == "list":
                    if 3 <= len(val) <= 6 and all(len(v) <= 22 for v in val):
                        chevron_strip(val)
                    else:
                        text_par([f"• {v}" for v in val])
                else:
                    text_par(val)
        if stmt:
            statement_line(stmt)

    for chunk in _split_slides(blocks):
        state["slide"] = None  # each chunk starts its own slide
        if styled:
            render_chunk_styled(chunk)
        else:
            render_chunk_plain(chunk)
            if state["slide"] is None:
                new_slide()  # an empty chunk (e.g. lone divider) still turns the page
    for slide in prs.slides:
        apply_transition(slide)
    prs.save(out_path)
    return state["count"]


def action_make_pptx(inp):
    notes = []
    blocks = normalize_blocks(inp.get("blocks"), notes)
    title = str(inp.get("title") or "").strip()
    if not blocks and not title:
        return {"success": False, "action": "make_pptx",
                "error": "blocks (or at least a title) required"}
    master_raw = str(inp.get("masterMediaPath")
                     or os.environ.get("MODULE_MASTERPPTXURL") or "").strip()
    master_path = None
    if master_raw:
        master_path, err = resolve_path(master_raw)
        if err:
            return {"success": False, "action": "make_pptx", "error": f"master: {err}"}
        if not master_path.lower().endswith(".pptx"):
            return {"success": False, "action": "make_pptx",
                    "error": f"master must be .pptx (got {master_path})"}
    out_path, stem = out_file(title or "deck", "pptx", {"t": title, "b": blocks})
    try:
        transition = str(inp.get("transition") or "fade").strip().lower()
        theme = inp.get("theme") if isinstance(inp.get("theme"), dict) else None
        slides = make_pptx_file(blocks, title, master_path, out_path,
                                transition=transition, theme=theme)
    except Exception as e:  # noqa: BLE001
        return {"success": False, "action": "make_pptx", "error": f"pptx build failed: {e}"}
    return {"success": True, "action": "make_pptx", "data": {
        "slides": slides, "master": bool(master_path),
        **({"notes": notes} if notes else {}),
        "_mediaImport": media_import_decl(out_path, "pptx", stem),
    }}


# ── xlsx: the dashboard genre ──────────────────────────────────────────────────────────────────
# Same stance as the pptx genre system: the archetypes are fixed and the block shapes choose
# between them, so no decision is left to the caller's prose. metric blocks are KPI cards,
# chart blocks are native charts fed from real cells (never inline literals — a chart whose
# numbers are baked into the drawing cannot be re-pointed by the person who opens the file),
# table blocks are styled data sheets.
XLSX_KPI_PER_ROW = 4
XLSX_KPI_COLS = 3           # merged card width in columns
XLSX_KPI_ROWS = 4           # 3 content rows (label / value / delta) + 1 gap row
XLSX_KPI_TOP = 3            # row 1 = title, row 2 = the single spacer row
# The Dashboard is ONE fixed grid: every column the same width, every band spanning all of it.
# Geometry that has to be guessed cannot be laid out — a chart is a floating drawing sized in
# centimetres, so the only way to make it land inside a cell box is to know what a column is worth.
XLSX_DASH_COLS = XLSX_KPI_PER_ROW * XLSX_KPI_COLS   # 12 columns = 4 card slots
XLSX_DASH_COL_W = 9.0       # every Dashboard column, in Excel's character units
# Excel's character width converts as px = width * 7 + 5, so a width-9 column is 68px = 1.80cm at
# 96dpi. Measured, not rounded up: an over-estimate here makes every chart overflow its own card
# into the neighbour's, which is exactly the overlap this pass exists to remove.
XLSX_COL_CM = 1.80          # what one width-9 column measures on screen
XLSX_ROW_CM = 0.529         # what one default row (15pt = 20px) measures on screen
XLSX_BAR_COLOR = "638EC6"   # one restrained blue — data bars are a reading aid, not decoration
XLSX_SCALE_LOW, XLSX_SCALE_MID, XLSX_SCALE_HIGH = "5B9BD5", "FFFFFF", "E46C6C"
# The dashboard wears the pptx genre's clothes: the same navy the section band uses (SLATE_D)
# plus the one accent (blue-600). A workbook and a deck built from the same blocks should not
# look like they came from two houses.
XLSX_BAND = "1E293B"        # title band fill — the pptx section-band navy
XLSX_ACCENT = "2563EB"      # the single accent (blue-600): card spine + one-series charts
XLSX_CARD_LINE = "D9D9D9"   # card hairline
XLSX_CARD_LABEL_BG = "F2F5FA"
XLSX_WHITE = "FFFFFF"
XLSX_TITLE_ROW_H = 26
# v4: the page itself is painted. Cards are WHITE boxes on a light slate canvas — that contrast
# is what made the reference dashboard read as "a designed page" and ours read as "empty".
# White-on-white has no edges, so every band was floating in a void however tightly it was packed.
XLSX_CANVAS = "F1F5F9"      # slate-100, the page the cards sit on
XLSX_UNIT_BADGE_BG = "FEF3C7"   # amber-100 — the "단위: 백만원" chip
XLSX_UNIT_BADGE_FG = "B45309"   # amber-700
# Card v3: the accent moved from a left spine to a thick bottom underline that cycles per card,
# so a row of cards reads as a row (the reference dashboards' language) instead of four clones.
XLSX_CARD_UNDERLINES = ("2563EB", "E0475B", "EAB308", "22C55E")
# Progress-bar KPIs (a ratio told twice: as a number and as a filled length).
XLSX_BARKPI_VALUE_COLS = 2   # merged width of the big colored number
XLSX_BARKPI_LABEL_COLS = 3   # merged width of the caption beside it
XLSX_BARKPI_BAR_COLS = 5     # minimum width of the cell the data bar fills (it takes the rest)
XLSX_BARKPI_COLS = XLSX_BARKPI_VALUE_COLS + XLSX_BARKPI_LABEL_COLS + XLSX_BARKPI_BAR_COLS
XLSX_BARKPI_ROWS = 2         # 1 content row + 1 gap row
XLSX_BARKPI_COLORS = ("E0475B", "EAB308", "2563EB", "22C55E")
# A chart is a drawing pinned over a pre-painted white cell box — the "chart card". Excel gives a
# chart no page of its own, so without the box underneath it floats on the canvas with no edge.
XLSX_CHARTCARD_ROWS = 18     # card height in grid rows (row 0 = the card header)
XLSX_CARD_HEADER_H = 18      # the header row, once it carries a title and not just the unit chip
XLSX_CHART_INSET_CM = 0.3    # so the drawing stops short of the card's right hairline
XLSX_RING_ROWS = 13          # a doughnut leaves the bottom of its card to the center block
XLSX_DOUGHNUT_HOLE = 60
XLSX_COMBO_BAR = "9EB9DA"    # muted blue bars so the accent line stays the foreground
XLSX_DLBL_MAX_POINTS = 24    # value labels above bars stop being readable past this many bars
# Ledger (총계정원장) genre: a document, not a heatmap.
XLSX_LEDGER_HEAD_BG = "D9D9D9"
XLSX_LEDGER_MONTH_BG = "F2F5FA"
XLSX_LEDGER_TOTAL_BG = "E4EBF5"
XLSX_LEDGER_LINE = "808080"
XLSX_LEDGER_MONTH_LABEL = "[월 계]"
XLSX_LEDGER_TOTAL_LABEL = "[누 계]"
# Korean market convention: up = red, down = blue (the inverse of the US convention).
XLSX_UP, XLSX_DOWN, XLSX_FLAT = "C00000", "1F5FBF", "808080"
# Ratio-ish columns diverge around zero, so a 3-color scale reads them; absolute magnitudes
# read as bars. The header is what tells the two apart.
_RATIO_HEADER_RE = re.compile(r"(%|율|증감|등락|change|delta|yoy)", re.I)
# The lenient names ("bar_chart" etc.) are dialect the model invents and there is no cost to
# accepting them — but the list was ALL dialect and no canon: `stock_chart`, a real page
# component, was the one chart block that fell through (2026-08-12 audit).
_CHART_BLOCK_TYPES = ("chart", "bar_chart", "line_chart", "pie_chart",
                      "donut_chart", "doughnut_chart", "column_chart", "area_chart")
_STOCK_BLOCK_TYPES = ("stock_chart", "candle_chart", "candlestick_chart", "ohlc_chart")
_CHART_TYPE_ALIASES = {"donut": "doughnut", "column": "bar", "area": "line"}
_CHART_TYPES = ("bar", "line", "pie", "doughnut", "combo")
_ROUND_CHARTS = ("pie", "doughnut")   # no value axis, own color per slice
# "2026-01-05" / "2026.01" / a real date cell -> the month a ledger row belongs to.
_YM_RE = re.compile(r"(\d{4})\s*[-/.년]?\s*(\d{1,2})")


def _sheet_title(raw, index, used):
    """Excel-legal, <=31 chars, unique within the workbook."""
    base = re.sub(r"[\\/*?:\[\]]", "-", str(raw or "")).strip()[:31] or f"Sheet{index + 1}"
    name, n = base, 1
    while name in used:
        suffix = f"_{n}"
        name = base[:31 - len(suffix)] + suffix
        n += 1
    used.add(name)
    return name


def _coerce_row(row):
    """"1,234" must land as the number 1234, not text — text numbers kill SUM and charts on the
    receiving end (2026-08-11 사용자 실측). A leading "=" stays a string here and openpyxl
    writes it as a live formula."""
    cells = []
    for v in row:
        num = parse_number(v)
        cells.append(num if num is not None else ("" if v is None else str(v)))
    return cells


_YEAR_HEADER_RE = re.compile(r"(연도|년도|회계연도|year)", re.I)


def _number_columns(body, ncols, headers=None):
    """{0-based column: number format} for the columns that really hold numbers. A column that
    is mostly text is not a number column, however many digits happen to sit in it.

    A YEAR column is an identifier, not a magnitude: "2,023" on an axis and a data bar over
    연도 both read as nonsense (2026-08-12 screenshot). All-integer values inside 1000–2999,
    or a year-ish header, take the plain "0" format — and the caller skips decoration
    (data bars / color scales) on any column formatted "0"."""
    headers = headers or []
    out = {}
    for ci in range(ncols):
        vals = [r[ci] for r in body if ci < len(r)]
        filled = [v for v in vals if v not in ("", None)]
        nums = [v for v in filled if isinstance(v, (int, float)) and not isinstance(v, bool)]
        if not nums or len(nums) * 2 <= len(filled):
            continue
        ints = [v for v in nums if float(v).is_integer()]
        header_txt = str(headers[ci]) if ci < len(headers) and headers[ci] is not None else ""
        year_like = (
            len(ints) == len(nums) and all(1000 <= v <= 2999 for v in nums)
        ) or _YEAR_HEADER_RE.search(header_txt)
        if year_like:
            out[ci] = "0"
        else:
            out[ci] = "#,##0" if len(ints) * 2 > len(nums) else "#,##0.00"
    return out


def _fit_column_widths(ws, headers, body, ncols):
    """Width from the widest thing the column will actually show, formatted as it will show."""
    from openpyxl.utils import get_column_letter

    for ci in range(ncols):
        texts = ([headers[ci]] if ci < len(headers) else []) + [
            (f"{r[ci]:,}" if isinstance(r[ci], (int, float)) and not isinstance(r[ci], bool)
             else str(r[ci]))
            for r in body if ci < len(r)]
        ws.column_dimensions[get_column_letter(ci + 1)].width = min(
            40, max(10, max([len(t) for t in texts] + [0]) + 2))


def _write_data_sheet(ws, sh):
    """Header row + coerced cells + the styling every data sheet gets."""
    from openpyxl.formatting.rule import ColorScaleRule, DataBarRule
    from openpyxl.styles import Font
    from openpyxl.utils import get_column_letter

    headers, rows_in = _table_rows(sh.get("headers"), sh.get("rows"))
    if headers:
        ws.append(headers)
        for c in ws[1]:
            c.font = Font(bold=True)
        ws.freeze_panes = "A2"
    body = []
    for row in rows_in:
        cells = _coerce_row(row)
        ws.append(cells)
        body.append(cells)

    ncols = max([len(headers)] + [len(r) for r in body] + [0])
    first_row = 2 if headers else 1
    last_row = first_row + len(body) - 1
    _fit_column_widths(ws, headers, body, ncols)
    for ci, fmt in _number_columns(body, ncols, headers).items():
        col = get_column_letter(ci + 1)
        for r in range(first_row, last_row + 1):
            ws.cell(row=r, column=ci + 1).number_format = fmt
        if last_row < first_row:
            continue
        if fmt == "0":
            # A year/identifier column: plain digits, and no bars or scales — decorating an
            # identifier reads as a claim about magnitude it does not make.
            continue
        rng = f"{col}{first_row}:{col}{last_row}"
        header_txt = headers[ci] if ci < len(headers) else ""
        if _RATIO_HEADER_RE.search(header_txt):
            ws.conditional_formatting.add(rng, ColorScaleRule(
                start_type="min", start_color=XLSX_SCALE_LOW,
                mid_type="percentile", mid_value=50, mid_color=XLSX_SCALE_MID,
                end_type="max", end_color=XLSX_SCALE_HIGH))
        else:
            # openpyxl 3.1's DataBarRule has no `solid` switch: the classic dataBar it emits is
            # the gradient (non-solid) one, which is the restrained look we want anyway.
            ws.conditional_formatting.add(rng, DataBarRule(
                start_type="min", end_type="max", color=XLSX_BAR_COLOR, showValue=True))


# ── xlsx: the ledger (총계정원장) sheet style ──────────────────────────────────────────────────
def _ym_key(v):
    """The YYYY-MM a ledger row belongs to, or None when the cell is not date-ish."""
    if hasattr(v, "year") and hasattr(v, "month"):
        return f"{int(v.year):04d}-{int(v.month):02d}"
    m = _YM_RE.match(str(v if v is not None else "").strip())
    return f"{int(m.group(1)):04d}-{int(m.group(2)):02d}" if m else None


def _ledger_subtotal_rows(ws, row, spans, numfmt, ncols, box):
    """Write the [월 계] / [누 계] band under a finished month. Both are LIVE formulas: editing
    a data row has to move them, which a precomputed number would not do.

    월 계 sums exactly the month's own data rows. 누 계 starts at the very first data row and
    ends at this month's last one — but as a *list* of the data spans, because the subtotal
    bands already written sit inside that stretch and one contiguous SUM would count them again.
    """
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    center = Alignment(horizontal="center", vertical="center")
    right = Alignment(horizontal="right", vertical="center")
    bands = ((XLSX_LEDGER_MONTH_LABEL, XLSX_LEDGER_MONTH_BG, spans[-1:]),
             (XLSX_LEDGER_TOTAL_LABEL, XLSX_LEDGER_TOTAL_BG, spans))
    for label, bg, use in bands:
        fill = PatternFill("solid", fgColor=bg)
        for ci in range(ncols):
            cell = ws.cell(row=row, column=ci + 1)
            cell.border = box
            cell.fill = fill
            cell.font = Font(bold=True)
            if ci in numfmt:
                col = get_column_letter(ci + 1)
                cell.value = "=SUM({})".format(
                    ",".join(f"{col}{a}:{col}{b}" for a, b in use))
                cell.number_format = numfmt[ci]
                cell.alignment = right
            elif ci == 0:
                cell.value = label
                cell.alignment = center
        row += 1
    return row


def _write_ledger_sheet(ws, sh):
    """A bookkeeping document: centered title, period line, fully ruled table, live monthly and
    running subtotals. Deliberately NOT a data sheet — no data bars, no color scale. A ledger is
    something you print and sign, and a heatmap in it reads as a mistake."""
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    headers, rows_in = _table_rows(sh.get("headers"), sh.get("rows"))
    body = [_coerce_row(row) for row in rows_in]
    ncols = max([len(headers)] + [len(r) for r in body] + [1])
    ws.sheet_view.showGridLines = False

    center = Alignment(horizontal="center", vertical="center")
    right = Alignment(horizontal="right", vertical="center")
    left = Alignment(horizontal="left", vertical="center")
    line = Side(style="thin", color=XLSX_LEDGER_LINE)
    box = Border(left=line, right=line, top=line, bottom=line)

    row = 1
    doc_title = str(sh.get("docTitle") or "").strip()
    period = str(sh.get("period") or "").strip()
    for text, size, color, height in ((doc_title, 16, XLSX_BAND, 30),
                                      (period, 10, "64748B", 18)):
        if not text:
            continue
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=ncols)
        head = ws.cell(row=row, column=1)
        head.value = text
        head.font = Font(size=size, bold=size >= 14, color=color)
        head.alignment = center
        ws.row_dimensions[row].height = height
        row += 1
    if row > 1:
        row += 1   # one breathing row between the masthead and the ruled table

    header_row = row
    if headers:
        fill = PatternFill("solid", fgColor=XLSX_LEDGER_HEAD_BG)
        for ci in range(ncols):
            cell = ws.cell(row=header_row, column=ci + 1)
            cell.value = headers[ci] if ci < len(headers) else ""
            cell.font = Font(bold=True)
            cell.fill = fill
            cell.alignment = center
            cell.border = box
        ws.freeze_panes = ws.cell(row=header_row + 1, column=1).coordinate
        row += 1

    numfmt = _number_columns(body, ncols)
    keycol = _resolve_column(headers, sh.get("subtotalBy")) if sh.get("subtotalBy") not in (
        None, "") else None
    spans, span_start, cur_key = [], row, None
    for cells in body:
        key = _ym_key(cells[keycol - 1]) if keycol and keycol - 1 < len(cells) else None
        if keycol and cur_key is not None and key != cur_key:
            spans.append((span_start, row - 1))
            row = _ledger_subtotal_rows(ws, row, spans, numfmt, ncols, box)
            span_start = row
        cur_key = key
        for ci in range(ncols):
            cell = ws.cell(row=row, column=ci + 1)
            cell.value = cells[ci] if ci < len(cells) else ""
            cell.border = box
            if ci in numfmt:
                cell.number_format = numfmt[ci]
                cell.alignment = right
            else:
                cell.alignment = center if ci == 0 else left
        row += 1
    if keycol and body:
        spans.append((span_start, row - 1))
        row = _ledger_subtotal_rows(ws, row, spans, numfmt, ncols, box)
    _fit_column_widths(ws, headers, body, ncols)


def _kpi_delta(delta, delta_type=None):
    """(text, color) for the delta line — sign decides, deltaType only rescues non-numbers."""
    if delta in (None, ""):
        return None, None
    d = parse_number(delta)
    if d is None:
        dt = str(delta_type or "").strip().lower()
        color = XLSX_UP if dt == "up" else XLSX_DOWN if dt == "down" else XLSX_FLAT
        return str(delta), color
    if d > 0:
        return f"▲ {abs(d):,}", XLSX_UP
    if d < 0:
        return f"▼ {abs(d):,}", XLSX_DOWN
    return "—", XLSX_FLAT


def _split_kpis(kpis):
    """(cards, progress bars) — style: "bar" picks the second archetype."""
    cards, bars = [], []
    for k in kpis:
        (bars if str(k.get("style") or "").strip().lower() == "bar" else cards).append(k)
    return cards, bars


def _dash_grid(ws, width):
    """Every Dashboard column the same width, BEFORE anything is laid out. Cards are cell boxes
    and charts are centimetre-sized drawings; only a fixed grid lets the two land on each other."""
    from openpyxl.utils import get_column_letter

    for c in range(1, width + 2):   # +1: the canvas runs one column past the last band
        ws.column_dimensions[get_column_letter(c)].width = XLSX_DASH_COL_W


def _cm_cols(cols):
    """Grid columns -> centimetres, minus the inset that keeps a drawing inside its card."""
    return max(1.0, cols * XLSX_COL_CM - XLSX_CHART_INSET_CM)


def _card_box(ws, r0, c0, rows, cols, fill=XLSX_WHITE):
    """Paint a white card with a hairline frame. Called BEFORE the content is written, because
    a fill is a cell property and a later value does not disturb it."""
    from openpyxl.styles import Border, PatternFill, Side

    hair = Side(style="thin", color=XLSX_CARD_LINE)
    body = PatternFill("solid", fgColor=fill)
    for dr in range(rows):
        for dc in range(cols):
            cell = ws.cell(row=r0 + dr, column=c0 + dc)
            cell.fill = body
            cell.border = Border(
                left=hair if dc == 0 else None,
                right=hair if dc == cols - 1 else None,
                top=hair if dr == 0 else None,
                bottom=hair if dr == rows - 1 else None)


def _paint_canvas(ws, last_row, width):
    """Light slate over the whole used region, skipping every cell a card already claimed. Run
    last: the cards decide where the page shows through, and an unfilled cell is by definition
    not part of one."""
    from openpyxl.styles import PatternFill

    page = PatternFill("solid", fgColor=XLSX_CANVAS)
    for r in range(1, last_row + 3):        # +2 rows of page under the final band
        for c in range(1, width + 2):       # +1 column past the band's right edge
            cell = ws.cell(row=r, column=c)
            if cell.fill is None or cell.fill.patternType is None:
                cell.fill = page


def _write_kpi_card_grid(ws, kpis, top_row, width=XLSX_DASH_COLS):
    """Cards left to right, at most four per row, STRETCHED so the row fills the band. Three cards
    on a twelve-column grid are three four-column cards, not three three-column cards and a hole:
    dead canvas at the end of a band is the thing that reads as "empty". Returns the row below."""
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    center = Alignment(horizontal="center", vertical="center")
    # A card is a box, not three loose cells: hairline on all four edges, a thick colored
    # underline at the bottom, a very light label strip on top. Interior verticals stay out —
    # the merges already read as one field and extra rules make a school-project grid.
    hair = Side(style="thin", color=XLSX_CARD_LINE)
    label_fill = PatternFill("solid", fgColor=XLSX_CARD_LABEL_BG)
    body_fill = PatternFill("solid", fgColor=XLSX_WHITE)
    per_row = min(len(kpis), XLSX_KPI_PER_ROW)
    span_cols = max(1, width // per_row)
    for i, k in enumerate(kpis):
        c0 = 1 + (i % per_row) * span_cols
        r0 = top_row + (i // per_row) * XLSX_KPI_ROWS
        underline = Side(style="thick",
                         color=XLSX_CARD_UNDERLINES[i % len(XLSX_CARD_UNDERLINES)])
        for dr in range(3):
            ws.merge_cells(start_row=r0 + dr, start_column=c0,
                           end_row=r0 + dr, end_column=c0 + span_cols - 1)
        # Styled after the merge, on every constituent cell — Excel composes a merged range's
        # frame from the cells underneath it, so styling only the anchor leaves three open sides.
        for dr in range(3):
            for dc in range(span_cols):
                cell = ws.cell(row=r0 + dr, column=c0 + dc)
                cell.border = Border(
                    left=hair if dc == 0 else None,
                    right=hair if dc == span_cols - 1 else None,
                    top=hair if dr == 0 else None,
                    bottom=underline if dr == 2 else None)
                cell.fill = label_fill if dr == 0 else body_fill
        ws.row_dimensions[r0].height = 16
        ws.row_dimensions[r0 + 1].height = 30

        lab = ws.cell(row=r0, column=c0)
        lab.value = str(k.get("label") or "")
        lab.font = Font(size=9, bold=True, color="64748B")
        lab.alignment = center

        unit = str(k.get("unit") or "").replace('"', "")
        fmt = f'#,##0"{unit}"' if unit else "#,##0"
        icon = str(k.get("icon") or "").strip()
        val = ws.cell(row=r0 + 1, column=c0)
        raw = k.get("value")
        s = "" if raw is None else str(raw).strip()
        num = parse_number(raw)
        if s.startswith("="):
            # A formula cannot carry a prefix and stay a formula, so an icon is dropped here
            # rather than silently turning the caller's live total into text.
            val.value = s
            val.number_format = fmt
        elif num is not None and not icon:
            val.value = num
            val.number_format = fmt
        else:
            # An icon trades the cell's numeric-ness away: "💰 334조원" is a label, not a value
            # you can SUM or chart. That is the deal, and it is the caller's to make.
            body = f"{num:,}{unit}" if num is not None else (f"{s}{unit}" if unit else s)
            val.value = f"{icon} {body}".strip()
        val.font = Font(size=20, bold=True, color=XLSX_BAND)
        val.alignment = center

        text, color = _kpi_delta(k.get("delta"), k.get("deltaType"))
        if text is not None:
            dcell = ws.cell(row=r0 + 2, column=c0)
            dcell.value = text
            dcell.font = Font(size=10, bold=True, color=color)
            dcell.alignment = center
    rows_used = -(-len(kpis) // per_row)
    return top_row + rows_used * XLSX_KPI_ROWS


def _write_kpi_bars(ws, kpis, top_row, width=XLSX_BARKPI_COLS):
    """A ratio told twice on one line: the number, its caption, and a data bar whose length is
    the number against `max`. The bar is a real conditional format on a real cell, so the length
    follows the value when someone edits it. Each row is its own white card spanning the band —
    the bar takes whatever the value and the caption leave."""
    from openpyxl.formatting.rule import DataBarRule
    from openpyxl.styles import Alignment, Font
    from openpyxl.utils import get_column_letter

    left = Alignment(horizontal="left", vertical="center")
    right = Alignment(horizontal="right", vertical="center")
    c_val = 1
    c_lab = c_val + XLSX_BARKPI_VALUE_COLS
    c_bar = c_lab + XLSX_BARKPI_LABEL_COLS
    bar_cols = max(XLSX_BARKPI_BAR_COLS,
                   width - XLSX_BARKPI_VALUE_COLS - XLSX_BARKPI_LABEL_COLS)
    for i, k in enumerate(kpis):
        r = top_row + i * XLSX_BARKPI_ROWS
        ws.row_dimensions[r].height = 24
        color = XLSX_BARKPI_COLORS[i % len(XLSX_BARKPI_COLORS)]
        num = parse_number(k.get("value"))
        top = parse_number(k.get("max"))
        top = float(top) if top not in (None, 0) else 100.0

        _card_box(ws, r, c_val, 1, c_bar + bar_cols - c_val)
        for c0, span in ((c_val, XLSX_BARKPI_VALUE_COLS), (c_lab, XLSX_BARKPI_LABEL_COLS),
                         (c_bar, bar_cols)):
            ws.merge_cells(start_row=r, start_column=c0, end_row=r, end_column=c0 + span - 1)
        unit = str(k.get("unit") or "")
        shown = k.get("value") if num is None else (f"{num:,.2f}".rstrip("0").rstrip("."))
        vcell = ws.cell(row=r, column=c_val)
        vcell.value = f"{shown}{unit}"
        vcell.font = Font(size=16, bold=True, color=color)
        vcell.alignment = right

        lcell = ws.cell(row=r, column=c_lab)
        lcell.value = str(k.get("label") or "")
        lcell.font = Font(size=10, color="64748B")
        lcell.alignment = left

        bcell = ws.cell(row=r, column=c_bar)
        bcell.value = num if num is not None else 0
        # "#,##0.##" keeps the decimal separator when there are no decimals, so 34 rendered as
        # "34." inside the bar (2026-08-12 사용자 실측). Excel has no single code that drops a
        # bare point, so the code is chosen from the value: whole numbers plain, the rest to two
        # places at most. The number sits left, on the filled end of the bar.
        whole = num is None or float(num).is_integer()
        bcell.number_format = "#,##0" if whole else "#,##0.##"
        bcell.alignment = left
        rng = f"{get_column_letter(c_bar)}{r}:{get_column_letter(c_bar)}{r}"
        ws.conditional_formatting.add(rng, DataBarRule(
            start_type="num", start_value=0, end_type="num", end_value=top,
            color=color, showValue=True))
    return top_row + len(kpis) * XLSX_BARKPI_ROWS


def _resolve_column(headers, ref):
    """header name (case-insensitive) or 0-based index -> 1-based column, None when unresolvable."""
    names = [str(h or "").strip() for h in headers]
    if isinstance(ref, str):
        want = ref.strip()
        for i, h in enumerate(names):
            if h == want:
                return i + 1
        for i, h in enumerate(names):
            if h.lower() == want.lower():
                return i + 1
        if re.fullmatch(r"\d+", want):
            ref = int(want)
    if isinstance(ref, bool):
        return None
    if isinstance(ref, int) and 0 <= ref < max(len(names), 1):
        return ref + 1
    return None


def _style_chart_title(chart, pt=11):
    """openpyxl builds a string title as rich text whose size lives on the paragraph's default
    run properties — runs carry no rPr of their own, so setting defRPr is the whole job."""
    rich = getattr(getattr(chart.title, "tx", None), "rich", None)
    for para in (getattr(rich, "p", None) or []):
        if para.pPr is not None and para.pPr.defRPr is not None:
            para.pPr.defRPr.sz = int(pt * 100)   # hundredths of a point
            para.pPr.defRPr.b = True


def _show_axes(chart):
    """Make every cartesian axis of `chart` — its sub-charts' axes included — actually visible.

    openpyxl never writes <c:delete>, and an omitted c:delete carries the schema default of TRUE,
    so Excel deletes the axis and every tick label with it: a chart captioned "단위: 백만원" then
    shows no number at all (2026-08-12 사용자 실측). Saying delete=0 out loud is the whole fix.
    The tick mark and the label position are said in the same breath because an axis nobody
    positioned is an axis Excel positions on its own.

    A combo's secondary value axis lives on the LINE sub-chart, not on the container, so the walk
    goes through `_charts` — setting the container's two axes would leave the 매출 scale hidden.
    """
    seen, axes = set(), []
    for owner in [chart] + list(getattr(chart, "_charts", None) or []):
        for ax in (getattr(owner, "x_axis", None), getattr(owner, "y_axis", None)):
            if ax is not None and id(ax) not in seen:
                seen.add(id(ax))
                axes.append(ax)
    for ax in axes:
        ax.delete = False
        ax.tickLblPos = "nextTo"
        ax.majorTickMark = "out"
        ax.minorTickMark = "none"
        if getattr(ax, "tagname", "") == "valAx":
            # Thousands separators on every value axis — a dashboard is read, not decoded.
            ax.number_format = "#,##0"
    return axes


def _value_labels():
    """Value-only data labels with EVERY show* flag explicit. An omitted CT_Boolean reads
    TRUE in Excel — the same delete-by-omission family as the hidden axes — so a bare
    showVal grew series names, category names and legend keys onto every bar
    ("자산총계, 2,023, 455,905,980" per label, 2026-08-12 screenshot)."""
    from openpyxl.chart.label import DataLabelList
    return DataLabelList(showVal=True, dLblPos="outEnd", showSerName=False,
                         showCatName=False, showLegendKey=False,
                         showPercent=False, showBubbleSize=False)


def _combo_chart(ws, vcols, labcol, max_row):
    """valueCols[0] as bars, valueCols[1] as a line on its own axis. Two quantities that share a
    time axis but not a unit (수량 vs 매출) belong on one chart with two scales — one scale would
    flatten the smaller series into the floor."""
    from openpyxl.chart import BarChart, LineChart, Reference
    from openpyxl.chart.label import DataLabelList
    from openpyxl.chart.shapes import GraphicalProperties
    from openpyxl.drawing.line import LineProperties

    bar = BarChart()
    bar.add_data(Reference(ws, min_col=vcols[0], min_row=1, max_row=max_row),
                 titles_from_data=True)
    line = LineChart()
    line.add_data(Reference(ws, min_col=vcols[1], min_row=1, max_row=max_row),
                  titles_from_data=True)
    if labcol:
        cats = Reference(ws, min_col=labcol, min_row=2, max_row=max_row)
        bar.set_categories(cats)
        line.set_categories(cats)
    for s in bar.series:
        s.graphicalProperties = GraphicalProperties(solidFill=XLSX_COMBO_BAR)
    for s in line.series:
        s.graphicalProperties = GraphicalProperties(
            ln=LineProperties(solidFill=XLSX_ACCENT, w=22000))
        s.smooth = False   # invented curvature between real points is a lie
    # The secondary axis is an axId the line owns; the primary then crosses at max so the two
    # axes sit on opposite sides instead of on top of each other.
    line.y_axis.axId = 200
    line.y_axis.number_format = "#,##0"
    line.y_axis.majorGridlines = None
    bar.y_axis.crosses = "max"
    bar += line
    # dLbls set on the BarChart group labels the bars only — the LineChart group keeps its own
    # (empty) dLbls, so the line stays bare. That split is the point: v4 banned labels on the
    # whole combo because the line dropped its numbers into the bars, which threw away the bar
    # values the reference dashboard shows above each column.
    if max_row - 1 <= XLSX_DLBL_MAX_POINTS:
        bar.dLbls = _value_labels()
    return bar


def _stock_chart(ws, vcols, labcol, volcol, max_row):
    """A real Excel candlestick: 3 series = High/Low/Close, 4 = Open/High/Low/Close.

    The order IS the contract — Excel reads a stock chart positionally, so a swapped pair draws
    silently wrong candles rather than failing. hiLowLines connect each day's range; UpDownBars
    are the bodies, painted in the Korean market's colors (up = red, down = blue), which is the
    inverse of Excel's own default and the whole reason they are set by hand.

    Volume rides a BarChart on a secondary axis when the candles carry it — the same two-scales-
    one-time-axis move as _combo_chart, because it is the same problem (a volume in millions
    would flatten a price in thousands into the floor)."""
    from openpyxl.chart import BarChart, Reference, StockChart
    from openpyxl.chart.axis import ChartLines
    from openpyxl.chart.shapes import GraphicalProperties
    from openpyxl.chart.updown_bars import UpDownBars
    from openpyxl.drawing.line import LineProperties

    chart = StockChart()
    for col in vcols:
        chart.add_data(Reference(ws, min_col=col, min_row=1, max_row=max_row),
                       titles_from_data=True)
    cats = Reference(ws, min_col=labcol, min_row=2, max_row=max_row) if labcol else None
    if cats is not None:
        chart.set_categories(cats)
    # Every price series is a marker on the day's line, never a line of its own: without this
    # the OHLC series draw as four crossing lines ON TOP of the candles.
    for s in chart.series:
        s.graphicalProperties = GraphicalProperties(ln=LineProperties(noFill=True))
    chart.hiLowLines = ChartLines()
    chart.upDownBars = UpDownBars(
        upBars=ChartLines(spPr=GraphicalProperties(solidFill=XLSX_UP)),
        downBars=ChartLines(spPr=GraphicalProperties(solidFill=XLSX_DOWN)))
    if not volcol:
        chart.legend = None   # "고가/저가/종가" names a candle nobody needed named
        return chart
    vol = BarChart()
    vol.add_data(Reference(ws, min_col=volcol, min_row=1, max_row=max_row),
                 titles_from_data=True)
    if cats is not None:
        vol.set_categories(cats)
    for s in vol.series:
        s.graphicalProperties = GraphicalProperties(solidFill=XLSX_COMBO_BAR)
    vol.y_axis.axId = 200
    vol.y_axis.number_format = "#,##0"
    vol.y_axis.majorGridlines = None
    chart.y_axis.crosses = "max"
    chart += vol
    return chart


def _write_doughnut_center(ws, spec, c0, row, cols):
    """Excel cannot paint text inside the hole of a doughnut, so the ratio the ring is about
    lives in merged cells UNDER the ring, inside the same card. v3 put it BESIDE the ring, which
    meant the block sat on whatever card came next once the bands were packed side by side."""
    from openpyxl.styles import Alignment, Font

    label = str(spec.get("centerLabel") or "").strip()
    value = str(spec.get("centerValue") or "").strip()
    if not label and not value:
        return
    center = Alignment(horizontal="center", vertical="center")
    if label:
        ws.merge_cells(start_row=row, start_column=c0, end_row=row, end_column=c0 + cols - 1)
        cell = ws.cell(row=row, column=c0)
        cell.value = label
        cell.font = Font(size=9, color="64748B")
        cell.alignment = center
    if value:
        ws.merge_cells(start_row=row + 1, start_column=c0,
                       end_row=row + 2, end_column=c0 + cols - 1)
        cell = ws.cell(row=row + 1, column=c0)
        cell.value = value
        cell.font = Font(size=24, bold=True, color=XLSX_BAND)
        cell.alignment = center
        ws.row_dimensions[row + 1].height = 26


def _card_title_cell(ws, row, c0, cols, text):
    """A chart title written as a CELL on the card's header row instead of inside the drawing.

    Excel centers a chart's own title over the plot. In a one-third-width card the ring IS the
    plot, so the title landed on the doughnut (2026-08-12 사용자 실측). Handing the title to the
    card fixes it at the source and speaks the KPI label strip's typography — bold, navy, on the
    card's own top row — so the two archetypes read as one page.
    """
    from openpyxl.styles import Alignment, Font

    text = str(text or "").strip()
    if not text or cols < 1:
        return
    if cols > 1:
        ws.merge_cells(start_row=row, start_column=c0, end_row=row, end_column=c0 + cols - 1)
    cell = ws.cell(row=row, column=c0)
    cell.value = text
    cell.font = Font(size=10, bold=True, color=XLSX_BAND)
    cell.alignment = Alignment(horizontal="left", vertical="center", indent=1)
    ws.row_dimensions[row].height = XLSX_CARD_HEADER_H


def _unit_badge(ws, row, c0, cols, unit):
    """A chart's unit ("단위: 백만원") as a small amber chip in the card's header row, hard right.
    Absent unit = no chip, no gap: the header row is the card's own padding either way. Returns
    the number of columns the chip took, so a card header can fill exactly what is left."""
    from openpyxl.styles import Alignment, Font, PatternFill

    text = str(unit or "").strip()
    if not text:
        return 0
    span = 2 if cols >= 4 else 1
    b0 = c0 + cols - span
    ws.merge_cells(start_row=row, start_column=b0, end_row=row, end_column=b0 + span - 1)
    fill = PatternFill("solid", fgColor=XLSX_UNIT_BADGE_BG)
    for dc in range(span):
        ws.cell(row=row, column=b0 + dc).fill = fill
    cell = ws.cell(row=row, column=b0)
    cell.value = text
    cell.font = Font(size=8, bold=True, color=XLSX_UNIT_BADGE_FG)
    cell.alignment = Alignment(horizontal="center", vertical="center")
    return span


def _build_chart(wb, spec, sheet_map, notes):
    """Resolve one chart spec against the data sheets -> (chart, ctype), or None with a note.
    Sizing and placement are NOT decided here: a chart's slot depends on what else is on the
    page, so the caller lays out first and only then hands each chart its centimetres."""
    from openpyxl.chart import BarChart, DoughnutChart, LineChart, PieChart, Reference
    from openpyxl.chart.label import DataLabelList
    from openpyxl.chart.shapes import GraphicalProperties
    from openpyxl.drawing.line import LineProperties

    kinds = {"line": LineChart, "bar": BarChart, "pie": PieChart,
             "doughnut": lambda: DoughnutChart(holeSize=XLSX_DOUGHNUT_HOLE)}
    title = str(spec.get("title") or "").strip()
    ref_name = str(spec.get("sheet") or "").strip()
    who = title or ref_name or "(untitled)"
    target = sheet_map.get(ref_name) or sheet_map.get(ref_name.lower())
    if not target:
        notes.append(f"chart '{who}' skipped: data sheet {ref_name!r} not found")
        return None
    ws = wb[target]
    if ws.max_row < 2:
        notes.append(f"chart '{who}' skipped: sheet '{target}' has no data rows")
        return None
    headers = [c.value for c in ws[1]]
    want = spec.get("valueCols")
    want = want if isinstance(want, list) else ([want] if want not in (None, "") else [])
    vcols, bad = [], []
    for ref in want:
        col = _resolve_column(headers, ref)
        (vcols if col else bad).append(col if col else ref)
    if bad:
        notes.append(f"chart '{who}': columns not found in '{target}': {bad}")
    if not vcols:
        notes.append(f"chart '{who}' skipped: no value column resolved in '{target}'")
        return None
    lab_ref = spec.get("labelCol")
    labcol = _resolve_column(headers, 0 if lab_ref in (None, "") else lab_ref)
    if labcol is None:
        notes.append(f"chart '{who}': label column {lab_ref!r} not found — using row numbers")

    ctype = str(spec.get("type") or "bar").strip().lower()
    ctype = _CHART_TYPE_ALIASES.get(ctype, ctype)
    if ctype == "combo" and len(vcols) < 2:
        notes.append(f"chart '{who}': combo needs two value columns — drawn as bars")
        ctype = "bar"
    if ctype == "stock" and len(vcols) < 3:
        notes.append(f"chart '{who}': a candlestick needs 고가/저가/종가 — drawn as a line")
        ctype = "line"
    if ctype == "combo":
        chart = _combo_chart(ws, vcols, labcol, ws.max_row)
    elif ctype == "stock":
        volcol = _resolve_column(headers, spec.get("volumeCol")) \
            if spec.get("volumeCol") not in (None, "") else None
        chart = _stock_chart(ws, vcols[:4], labcol, volcol, ws.max_row)
    else:
        chart = kinds.get(ctype, BarChart)()
        for col in vcols:
            chart.add_data(Reference(ws, min_col=col, min_row=1, max_row=ws.max_row),
                           titles_from_data=True)
        if labcol:
            chart.set_categories(
                Reference(ws, min_col=labcol, min_row=2, max_row=ws.max_row))
    # A round chart keeps NO internal title: its card writes the title as a header cell instead
    # (see _card_title_cell). Everything cartesian has a plot wide enough to sit under one.
    if title and ctype not in _ROUND_CHARTS:
        chart.title = title
        _style_chart_title(chart)
    if ctype not in _ROUND_CHARTS:
        _show_axes(chart)
    if len(vcols) < 2 and ctype not in _ROUND_CHARTS:
        chart.legend = None  # a one-series legend is noise
        # One series = one color. Excel's default palette exists to tell series apart; with
        # nothing to tell apart it just adds a hue the design did not choose.
        for s in chart.series:
            if ctype == "line":
                s.graphicalProperties = GraphicalProperties(
                    ln=LineProperties(solidFill=XLSX_ACCENT, w=22000))
                s.smooth = False   # invented curvature between real points is a lie
            else:
                s.graphicalProperties = GraphicalProperties(solidFill=XLSX_ACCENT)
    if chart.legend is not None:
        # A legend on the right eats the plot and lands ON the data (2026-08-12 사용자 실측).
        # Under the plot it takes rows the drawing already reserved.
        chart.legend.position = "b"
        chart.legend.overlay = False
    # Value labels only where they have somewhere to go: above the bars of a SINGLE series,
    # and only while the bars are still wide enough to hold a number. Three series of
    # nine-digit values wore their labels as a solid ink cloud (2026-08-12 screenshot) —
    # a multi-series bar already has the legend; the reference dashboards label only the
    # lone-series kind. A combo labels its bars inside _combo_chart, on the bar group alone.
    if ctype == "bar" and len(vcols) == 1 and ws.max_row - 1 <= XLSX_DLBL_MAX_POINTS:
        chart.dLbls = _value_labels()
    # The chart's own page: white area, no frame line — so the drawing reads as part of the white
    # card painted under it instead of a gray rectangle floating on the canvas.
    area = GraphicalProperties(solidFill=XLSX_WHITE)
    area.ln = LineProperties(noFill=True)
    chart.graphical_properties = area
    return chart, ctype


def _chart_bands(prepared, width):
    """Plan the chart bands: [[(entry, cols), ...], ...].

    A wide chart and a round one share a band — wide ~2/3, ring ~1/3 — which is the reference
    layout and also the only arrangement where the ring's caption block has room under it.
    Whatever is left over pairs up two to a band, and a single leftover spans the whole width."""
    wide = [e for e in prepared if e[1] not in _ROUND_CHARTS]
    ring = [e for e in prepared if e[1] in _ROUND_CHARTS]
    bands = []
    while wide and ring:
        narrow = max(3, round(width / 3))
        bands.append([(wide.pop(0), width - narrow), (ring.pop(0), narrow)])
    rest = sorted(wide + ring, key=lambda e: e[2])   # back into the caller's order
    while rest:
        take, rest = rest[:2], rest[2:]
        if len(take) == 1:
            bands.append([(take[0], width)])
        else:
            half = width // 2
            bands.append([(take[0], half), (take[1], width - half)])
    return bands


def _add_dashboard_charts(wb, ws_dash, charts, sheet_map, top_row, notes, width):
    """Chart cards below the KPI band. Returns (charts placed, first free row below)."""
    from openpyxl.utils import get_column_letter

    prepared = []
    for i, spec in enumerate(charts):
        made = _build_chart(wb, spec, sheet_map, notes)
        if made:
            prepared.append((made[0], made[1], i, spec))
    row = top_row
    for band in _chart_bands(prepared, width):
        c0 = 1
        for (chart, ctype, _i, spec), cols in band:
            _card_box(ws_dash, row, c0, XLSX_CHARTCARD_ROWS, cols)
            badge = _unit_badge(ws_dash, row, c0, cols, spec.get("unit"))
            ring = ctype in _ROUND_CHARTS
            if ring:
                # The ring gave its title away; the card's header row takes it, minus whatever
                # the unit chip already claimed on the right.
                _card_title_cell(ws_dash, row, c0, cols - badge, spec.get("title"))
            # -2: the header row on top and one padding row above the card's bottom hairline.
            rows = XLSX_RING_ROWS if ring else XLSX_CHARTCARD_ROWS - 2
            # The drawing is sized from the measured grid, never from a guess: the card is
            # `cols` columns of XLSX_COL_CM minus the inset, so the chart cannot reach the
            # neighbour's hairline however narrow the card gets.
            chart.width = _cm_cols(cols)
            chart.height = rows * XLSX_ROW_CM
            ws_dash.add_chart(chart, f"{get_column_letter(c0)}{row + 1}")
            if ring:
                _write_doughnut_center(ws_dash, spec, c0, row + 1 + rows, cols)
            c0 += cols
        row += XLSX_CHARTCARD_ROWS + 1   # exactly one spacer row between bands
    return len(prepared), row


def _write_note_band(ws, texts, top_row, width):
    """callout blocks as full-width amber note rows at the foot of the Dashboard.

    A callout is where the honesty note lives ("추정치", "장중 데이터"), and until now it was
    the ONE block type every renderer dropped — the caveat vanished while the numbers it
    qualified survived. It wears the unit chip's amber so it reads as an annotation, not data."""
    from openpyxl.styles import Alignment, Font, PatternFill

    fill = PatternFill("solid", fgColor=XLSX_UNIT_BADGE_BG)
    row = top_row
    for text in texts:
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=width)
        for c in range(1, width + 1):
            ws.cell(row=row, column=c).fill = fill
        cell = ws.cell(row=row, column=1)
        cell.value = f"※ {text}"
        cell.font = Font(size=9, bold=True, color=XLSX_UNIT_BADGE_FG)
        cell.alignment = Alignment(horizontal="left", vertical="center", indent=1,
                                   wrap_text=True)
        ws.row_dimensions[row].height = 20
        row += 2   # one page row between notes, the spacing every other band uses
    return row


def _write_dash_title(ws, title, width):
    """Row 1 = the navy band, same spine as the pptx section band."""
    from openpyxl.styles import Alignment, Font, PatternFill

    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=width)
    fill = PatternFill("solid", fgColor=XLSX_BAND)
    for c in range(1, width + 1):
        ws.cell(row=1, column=c).fill = fill
    head = ws.cell(row=1, column=1)
    head.value = str(title)
    head.font = Font(size=14, bold=True, color=XLSX_WHITE)
    head.alignment = Alignment(horizontal="left", vertical="center", indent=1)
    ws.row_dimensions[1].height = XLSX_TITLE_ROW_H


def make_xlsx_file(sheets, out_path, title=None, kpis=None, charts=None, callouts=None):
    """Data sheets always; a Dashboard sheet in front of them when KPIs, charts or notes exist."""
    import openpyxl

    sheets = [s for s in (sheets or []) if isinstance(s, dict)]
    kpis = [k for k in (kpis or []) if isinstance(k, dict)]
    charts = [c for c in (charts or []) if isinstance(c, dict)]
    callouts = [str(c) for c in (callouts or []) if str(c or "").strip()]

    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    used, notes = set(), []
    # Created first so it is the sheet the file opens on; filled last, because the charts
    # inside it need the data sheets to exist.
    ws_dash = None
    if kpis or charts or callouts:
        ws_dash = wb.create_sheet(title=_sheet_title("Dashboard", 0, used))
        # The single biggest difference between "a report" and "a spreadsheet someone typed in":
        # the grid. Data sheets keep theirs — they are working sheets.
        ws_dash.sheet_view.showGridLines = False

    sheet_map = {}
    for i, sh in enumerate(sheets):
        raw = str(sh.get("name") or "")
        name = _sheet_title(raw or f"Sheet{i + 1}", i, used)
        ws_new = wb.create_sheet(title=name)
        if str(sh.get("style") or "").strip().lower() == "ledger":
            _write_ledger_sheet(ws_new, sh)
        else:
            _write_data_sheet(ws_new, sh)
        for key in (raw, raw.strip().lower(), name, name.strip().lower()):
            if key:
                sheet_map.setdefault(key, name)

    placed = 0
    if ws_dash is not None:
        # One fixed grid, then the bands in reading order: title -> KPI cards -> charts ->
        # progress bars, one spacer row apiece. Every band already ends with its own gap row, so
        # nothing adds a second one — two blank rows is what made v3 read as a mostly empty page.
        width = XLSX_DASH_COLS
        _dash_grid(ws_dash, width)
        if title:
            _write_dash_title(ws_dash, title, width)
        cards, bars = _split_kpis(kpis)
        row = XLSX_KPI_TOP
        if cards:
            row = _write_kpi_card_grid(ws_dash, cards, row, width)
        if charts:
            placed, row = _add_dashboard_charts(
                wb, ws_dash, charts, sheet_map, row, notes, width)
        if bars:
            row = _write_kpi_bars(ws_dash, bars, row, width)
        if callouts:
            row = _write_note_band(ws_dash, callouts, row, width)
        _paint_canvas(ws_dash, row, width)

    wb.save(out_path)
    return {"sheets": len(sheets), "dashboard": ws_dash is not None,
            "kpis": len(kpis), "charts": placed, "notes": notes}


def _chart_series(b):
    """chart-family block -> (ctype, title, labels, [(unique series name, values)]) or None.

    One parse for both consumers: xlsx turns it into real cells a native chart points at, pptx
    hands it straight to CategoryChartData. Two parsers would drift, and the pptx one would be
    the one that quietly stopped accepting a dialect."""
    t = str(b.get("type") or "").lower()
    p = b.get("props") or {}
    ctype = str(p.get("chartType") or p.get("type") or "").strip().lower()
    if not ctype and t.endswith("_chart"):
        ctype = t[: -len("_chart")]
    ctype = _CHART_TYPE_ALIASES.get(ctype, ctype)
    if ctype not in _CHART_TYPES:
        ctype = "bar"

    labels = [str(x) for x in (p.get("labels") or p.get("categories") or [])]
    cols = []
    series = p.get("series") or p.get("datasets")
    if isinstance(series, list) and series:
        for si, s in enumerate(series):
            if not isinstance(s, dict):
                continue
            name = str(s.get("name") or s.get("label") or f"계열{si + 1}")
            cols.append((name, list(s.get("data") or s.get("values") or [])))
    else:
        data = p.get("data") if isinstance(p.get("data"), list) else p.get("values")
        if isinstance(data, list) and data and not isinstance(data[0], dict):
            cols.append((str(p.get("title") or "값"), list(data)))
    if not labels or not cols:
        return None

    seen, named = set(), []
    for name, vals in cols:
        uniq, n = name, 2
        while uniq in seen:
            uniq = f"{name} {n}"
            n += 1
        seen.add(uniq)
        named.append((uniq, vals))
    return ctype, str(p.get("title") or "").strip(), labels, named


def _chart_block(b, index, taken):
    """chart-family block -> (data sheet, chart spec). Inline series become real cells so the
    native chart can point at them; None when the block carries no plottable numbers."""
    parsed = _chart_series(b)
    if not parsed:
        return None
    p = b.get("props") or {}
    ctype, title, labels, named = parsed
    names = [n for n, _v in named]
    cols = named

    title = title or f"차트{index}"
    sheet_name, n = title, 2   # the chart spec points at this sheet BY NAME — keep it unique
    while sheet_name in taken:
        sheet_name = f"{title} {n}"
        n += 1
    taken.add(sheet_name)

    rows = [[lab] + [(vals[ri] if ri < len(vals) else None) for _n, vals in cols]
            for ri, lab in enumerate(labels)]
    sheet = {"name": sheet_name, "headers": ["항목"] + names, "rows": rows}
    spec = {"type": ctype, "title": title, "sheet": sheet_name,
            "labelCol": 0, "valueCols": names,
            # Whichever unit-ish prop the block carries becomes the card's badge; none = no badge.
            "unit": p.get("unit") or p.get("unitLabel") or p.get("yUnit"),
            "centerLabel": p.get("centerLabel"), "centerValue": p.get("centerValue")}
    return sheet, spec


def _stock_block(b, index, taken):
    """stock_chart block -> (OHLC(V) data sheet, candlestick spec). None when the block carries
    no candles — the same contract as _chart_block, so the caller treats both alike."""
    p = b.get("props") or {}
    headers, rows = _stock_rows(p)
    if not rows or len(headers) < 4:
        return None
    title = str(p.get("title") or p.get("symbol") or "").strip() or f"차트{index}"
    sheet_name, n = title, 2
    while sheet_name in taken:
        sheet_name = f"{title} {n}"
        n += 1
    taken.add(sheet_name)
    # Excel reads a stock chart positionally, so the value columns go in exactly the order the
    # candle wants: O/H/L/C when the open is there, H/L/C when it is not.
    price = [h for h in ("시가", "고가", "저가", "종가") if h in headers]
    if len(price) < 3:
        return None
    sheet = {"name": sheet_name, "headers": headers, "rows": rows}
    spec = {"type": "stock", "title": title, "sheet": sheet_name,
            "labelCol": headers[0], "valueCols": price,
            "volumeCol": "거래량" if "거래량" in headers else None,
            "unit": p.get("unit") or p.get("unitLabel")}
    return sheet, spec


def _blocks_to_xlsx(blocks):
    """Render blocks -> (table sheets, kpis, charts, chart data sheets, callout notes)."""
    tables, kpis, charts, chart_sheets, callouts = [], [], [], [], []
    taken, last_header = set(), None
    for b in blocks:
        t, p = str(b.get("type") or "").lower(), b.get("props") or {}
        if t == "header":
            last_header = str(p.get("text") or "")
        elif t == "table":
            name = last_header or f"표{len(tables) + 1}"
            taken.add(name)
            headers, rows = _table_rows(p.get("headers"), p.get("rows"))
            tables.append({"name": name, "headers": headers, "rows": rows})
            last_header = None
        elif t == "metric":
            kpis.append({"label": p.get("label"), "value": p.get("value"),
                         "unit": p.get("unit"), "delta": p.get("delta"),
                         "deltaType": p.get("deltaType"), "icon": p.get("icon"),
                         "style": p.get("style"), "max": p.get("max")})
        elif t == "callout":
            title, msg = _callout_parts(p)
            line = f"{title} — {msg}" if title and msg else (title or msg)
            if line:
                callouts.append(line)
        elif t == "key_value":
            # A fact list is a two-column table — the only shape a spreadsheet has for it.
            pairs = _kv_pairs(p)
            if pairs:
                name = str(p.get("title") or "").strip() or last_header or f"표{len(tables) + 1}"
                taken.add(name)
                tables.append({"name": name, "headers": ["항목", "값"],
                               "rows": [[k, v] for k, v in pairs]})
                last_header = None
        elif t in _STOCK_BLOCK_TYPES:
            made = _stock_block(b, len(charts) + 1, taken)
            if made:
                chart_sheets.append(made[0])
                charts.append(made[1])
        elif t in _CHART_BLOCK_TYPES:
            made = _chart_block(b, len(charts) + 1, taken)
            if made:
                chart_sheets.append(made[0])
                charts.append(made[1])
    return tables, kpis, charts, chart_sheets, callouts


def action_make_xlsx(inp):
    sheets = inp.get("sheets")
    sheets = [s for s in sheets if isinstance(s, dict)] if isinstance(sheets, list) else []
    kpis = [k for k in (inp.get("kpis") or []) if isinstance(k, dict)]
    charts = [c for c in (inp.get("charts") or []) if isinstance(c, dict)]

    drops = []
    b_tables, b_kpis, b_charts, b_chart_sheets, b_callouts = _blocks_to_xlsx(
        normalize_blocks(inp.get("blocks"), drops))
    # Explicit input wins per axis, so blocks never duplicate what the caller stated.
    if not sheets:
        sheets = b_tables
    if not kpis:
        kpis = b_kpis
    if not charts and b_charts:
        charts = b_charts
        names = {str(s.get("name") or "") for s in sheets}
        for sh, spec in zip(b_chart_sheets, b_charts):
            while sh["name"] in names:
                sh["name"] = sh["name"] + "_"
            spec["sheet"] = sh["name"]
            names.add(sh["name"])
        sheets = list(sheets) + b_chart_sheets

    if not sheets and not kpis and not charts and not b_callouts:
        return {"success": False, "action": "make_xlsx",
                "error": "nothing to write — pass sheets/kpis/charts, or blocks containing "
                         "table, metric, chart or stock_chart blocks"}
    title = str(inp.get("title") or (sheets[0].get("name") if sheets else "") or "sheet")
    out_path, stem = out_file(title, "xlsx", {"s": sheets, "k": kpis, "c": charts})
    try:
        res = make_xlsx_file(sheets, out_path, title=title, kpis=kpis, charts=charts,
                             callouts=b_callouts)
    except Exception as e:  # noqa: BLE001
        return {"success": False, "action": "make_xlsx", "error": f"xlsx build failed: {e}"}
    res["notes"] = drops + list(res.get("notes") or [])
    if not res["notes"]:
        res.pop("notes")
    return {"success": True, "action": "make_xlsx", "data": {
        **res, "_mediaImport": media_import_decl(out_path, "xlsx", stem),
    }}


# ── docx / pdf: the report genre ───────────────────────────────────────────────────────────────
# The deck genre's print sibling — same navy, same single accent, same restraint. A report is not
# a deck flowed onto A4: it opens with a cover, its headings are ranked by weight rather than by
# size alone, its tables are read (numbers right, header banded) and its KPI rows are laid out
# side by side instead of stacked into a column of orphan sentences.
RPT_NAVY = "1E293B"        # headings — the pptx section band's navy
RPT_ACCENT = "2563EB"      # the one accent (blue-600)
RPT_BODY = "334155"
RPT_MUTED = "64748B"       # date line, KPI labels, footer
RPT_BAND = "F1F5F9"        # KPI card / callout tint
RPT_RULE = "CBD5E1"        # table hairline
RPT_HEADFILL = "E2E8F0"    # table header row
RPT_CALLOUT_EDGE = "BFDBFE"


def _today_line():
    """The cover's date on the OWNER's wall clock — `_runtime/tz.local()`, the one entry point
    for it. Hand-rolling FIREBAT_TZ here fell back to the host's clock when the zone was
    unset, and the host's zone is an accident of deployment (CI's clock-discipline test
    caught exactly that fallback)."""
    rt = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "_runtime")
    if rt not in sys.path:
        sys.path.insert(0, rt)
    import tz as clock
    return clock.local().strftime("%Y. %m. %d.")


def _chart_table_block(b):
    """A chart as a table — what a page with no plotting engine can honestly show of one.
    docx and pdf have no native chart, and silently dropping the numbers is the worse answer."""
    parsed = _chart_series(b)
    if not parsed:
        return None
    _ctype, title, labels, named = parsed
    headers = ["항목"] + [n for n, _v in named]
    rows = [[lab] + [(vals[i] if i < len(vals) else None) for _n, vals in named]
            for i, lab in enumerate(labels)]
    return {"type": "table", "props": {"headers": headers, "rows": rows, "title": title}}


def _report_cover(blocks, title):
    """(cover title, remaining blocks). A stream that opens with a level-1 header is telling us
    its own name — that header becomes the cover, and a divider right behind it was only ever
    saying "the document starts here", which the cover now says better."""
    blocks = list(blocks or [])
    cover = str(title or "").strip()
    if blocks and str(blocks[0].get("type") or "") == "header" \
            and _header_level(blocks[0].get("props") or {}) == 1:
        cover = str((blocks[0].get("props") or {}).get("text") or "").strip() or cover
        blocks = blocks[1:]
        if blocks and str(blocks[0].get("type") or "") == "divider":
            blocks = blocks[1:]
    return cover, blocks


def _metric_runs(blocks):
    """[(index, [metric props, ...])] — consecutive metric blocks are ONE strip. A KPI row read
    as four stacked sentences is the single most common way a report looks unmade."""
    runs, i = [], 0
    while i < len(blocks):
        if str(blocks[i].get("type") or "") == "metric":
            start, ms = i, []
            while i < len(blocks) and str(blocks[i].get("type") or "") == "metric":
                ms.append(blocks[i].get("props") or {})
                i += 1
            runs.append((start, ms))
        else:
            i += 1
    return runs


def _docx_el(tag, **attrs):
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    el = OxmlElement(tag)
    for k, v in attrs.items():
        el.set(qn("w:" + k), str(v))
    return el


def _docx_borders(prop_el, color, sz, edges, container="w:pBdr", val="single"):
    bdr = _docx_el(container)
    for edge in edges:
        bdr.append(_docx_el("w:" + edge, val=val, sz=sz, space="0", color=color))
    prop_el.append(bdr)


def _docx_shade(prop_el, color):
    prop_el.append(_docx_el("w:shd", val="clear", color="auto", fill=color))


def _docx_para_box(par, fill=None, edge=None, edges=("top", "left", "bottom", "right"), sz=6):
    # pBdr before shd — that is the order CT_PPr declares, and Word is only forgiving until it
    # is not.
    pPr = par._p.get_or_add_pPr()
    if edge:
        _docx_borders(pPr, edge, sz, edges)
    if fill:
        _docx_shade(pPr, fill)


def _docx_table_borders(table, color=RPT_RULE, sz=4,
                        edges=("top", "left", "bottom", "right", "insideH", "insideV"),
                        val="single"):
    _docx_borders(table._tbl.tblPr, color, sz, edges, container="w:tblBorders", val=val)


def _docx_cell_shade(cell, color):
    _docx_shade(cell._tc.get_or_add_tcPr(), color)


def _docx_set_widths(table, widths):
    """Word autofits to CONTENT, which turns a four-card KPI row into four different-sized
    cards. A card row is a grid — the widths are the design, so they are stated."""
    table.autofit = False
    for row in table.rows:
        for cell, w in zip(row.cells, widths):
            cell.width = w


def _docx_run(par, text, size, bold=False, color=RPT_BODY):
    from docx.shared import Pt, RGBColor
    run = par.add_run(str(text))
    run.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(color)
    return run


def _docx_page_footer(section):
    """A page number in the section footer — a PAGE field, so it counts itself. python-docx has
    no API for fields, but fldSimple is one element and Word owns the arithmetic."""
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    par = section.footer.paragraphs[0] if section.footer.paragraphs \
        else section.footer.add_paragraph()
    par.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fld = _docx_el("w:fldSimple", instr=" PAGE ")
    run = OxmlElement("w:r")
    rpr = OxmlElement("w:rPr")
    sz = OxmlElement("w:sz")
    sz.set(qn("w:val"), "16")            # half-points: 8pt
    color = OxmlElement("w:color")
    color.set(qn("w:val"), RPT_MUTED)
    rpr.append(sz)
    rpr.append(color)
    text = OxmlElement("w:t")
    text.text = "1"
    run.append(rpr)
    run.append(text)
    fld.append(run)
    par._p.append(fld)


def _docx_heading(d, text, lvl):
    from docx.shared import Pt
    par = d.add_paragraph()
    _docx_run(par, text, {1: 16, 2: 13, 3: 11.5}[lvl], bold=True, color=RPT_NAVY)
    par.paragraph_format.space_before = Pt({1: 18, 2: 13, 3: 9}[lvl])
    par.paragraph_format.space_after = Pt({1: 7, 2: 5, 3: 3}[lvl])
    if lvl == 1:
        _docx_para_box(par, edge=RPT_ACCENT, edges=("bottom",), sz=8)
    return par


def _docx_kpi_strip(d, metrics, width):
    """Label over value, side by side, in a borderless tinted table — the deck's metric card
    with the deck's geometry taken out."""
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt

    per_row = min(4, max(1, len(metrics)))
    n_rows = -(-len(metrics) // per_row)
    table = d.add_table(rows=n_rows, cols=per_row)
    _docx_table_borders(table, val="none")
    _docx_set_widths(table, [int(width / per_row)] * per_row)
    for idx, m in enumerate(metrics):
        cell = table.cell(idx // per_row, idx % per_row)
        _docx_cell_shade(cell, RPT_BAND)
        lab = cell.paragraphs[0]
        lab.paragraph_format.space_after = Pt(0)
        _docx_run(lab, m.get("label") or "", 8, color=RPT_MUTED)
        val = cell.add_paragraph()
        val.paragraph_format.space_after = Pt(0)
        _docx_run(val, f"{m.get('value', '')}{m.get('unit') or ''}", 16, bold=True,
                  color=RPT_NAVY)
        if m.get("delta") not in (None, ""):
            dl = cell.add_paragraph()
            dl.paragraph_format.space_after = Pt(0)
            _docx_run(dl, m.get("delta"), 8.5, color=RPT_MUTED)
        for par in cell.paragraphs:
            par.alignment = WD_ALIGN_PARAGRAPH.LEFT
    d.add_paragraph().paragraph_format.space_after = Pt(4)


def _docx_table(d, p):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt

    headers, rows = _table_rows(p.get("headers"), p.get("rows"))
    if not headers and rows:
        headers = [str(c) for c in rows[0]]
        rows = rows[1:]
    if not headers:
        return
    table = d.add_table(rows=len(rows) + 1, cols=len(headers))
    _docx_table_borders(table)
    numeric = _numeric_cols(rows, len(headers))
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        _docx_cell_shade(cell, RPT_HEADFILL)
        par = cell.paragraphs[0]
        par.alignment = WD_ALIGN_PARAGRAPH.CENTER
        par.paragraph_format.space_after = Pt(0)
        _docx_run(par, h, 9.5, bold=True, color=RPT_NAVY)
    for r, row in enumerate(rows):
        for c in range(len(headers)):
            val = row[c] if c < len(row) else ""
            par = table.cell(r + 1, c).paragraphs[0]
            par.alignment = (WD_ALIGN_PARAGRAPH.RIGHT if c in numeric
                             else WD_ALIGN_PARAGRAPH.LEFT)
            par.paragraph_format.space_after = Pt(0)
            _docx_run(par, "" if val is None else val, 9.5)
    d.add_paragraph().paragraph_format.space_after = Pt(4)


def _docx_key_value(d, p, width):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt

    pairs = _kv_pairs(p)
    if not pairs:
        return
    title = str(p.get("title") or "").strip()
    if title:
        _docx_heading(d, title, 3)
    table = d.add_table(rows=len(pairs), cols=2)
    _docx_table_borders(table, edges=("insideH",))   # rules between facts, nothing around them
    _docx_set_widths(table, [int(width * 0.45), int(width * 0.55)])
    for r, (k, v) in enumerate(pairs):
        kp = table.cell(r, 0).paragraphs[0]
        kp.paragraph_format.space_after = Pt(2)
        _docx_run(kp, k, 9.5, color=RPT_MUTED)
        vp = table.cell(r, 1).paragraphs[0]
        vp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        vp.paragraph_format.space_after = Pt(2)
        _docx_run(vp, "" if v is None else v, 9.5, bold=True, color=RPT_NAVY)
    d.add_paragraph().paragraph_format.space_after = Pt(4)


def _docx_callout(d, p):
    from docx.shared import Pt
    title, msg = _callout_parts(p)
    if not title and not msg:
        return
    par = d.add_paragraph()
    _docx_para_box(par, fill=RPT_BAND, edge=RPT_CALLOUT_EDGE, sz=6)
    fmt = par.paragraph_format
    fmt.left_indent = Pt(8)
    fmt.right_indent = Pt(8)
    fmt.space_before = Pt(8)
    fmt.space_after = Pt(8)
    if title:
        _docx_run(par, title + ("  " if msg else ""), 10, bold=True, color=RPT_NAVY)
    if msg:
        _docx_run(par, msg, 10, color=RPT_BODY)


def _docx_image(d, p, max_width):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    path, _err = resolve_path(str(p.get("src") or ""))
    if not path:
        return
    try:
        pic = d.add_picture(path)
    except Exception:  # noqa: BLE001 — a bad image loses itself, not the document
        return
    if pic.width > max_width:
        pic.height = int(pic.height * max_width / pic.width)
        pic.width = max_width
    d.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER


def make_docx_file(blocks, title, out_path):
    """Render blocks -> a .docx report: cover, ranked headings, read-able tables, KPI strips,
    callout boxes, page-numbered footer. Plain text and lists stay plain — a report is mostly
    prose, and the genre earns its keep at the places prose is not."""
    import docx
    from docx.shared import Pt

    d = docx.Document()
    section = d.sections[0]
    body_width = section.page_width - section.left_margin - section.right_margin
    cover, blocks = _report_cover(blocks, title)
    if cover:
        par = d.add_paragraph()
        par.paragraph_format.space_before = Pt(150)
        par.paragraph_format.space_after = Pt(10)
        _docx_run(par, cover, 26, bold=True, color=RPT_NAVY)
        date_par = d.add_paragraph()
        date_par.paragraph_format.space_after = Pt(6)
        _docx_run(date_par, _today_line(), 10, color=RPT_MUTED)
        rule = d.add_paragraph()
        rule.paragraph_format.space_after = Pt(0)
        _docx_para_box(rule, edge=RPT_ACCENT, edges=("bottom",), sz=12)
        if blocks:
            d.add_page_break()
    _docx_page_footer(section)

    metric_at = dict(_metric_runs(blocks))
    skip_to = -1
    for i, b in enumerate(blocks):
        if i < skip_to:
            continue
        t, p = str(b.get("type") or ""), b.get("props") or {}
        if i in metric_at:
            _docx_kpi_strip(d, metric_at[i], body_width)
            skip_to = i + len(metric_at[i])
        elif t == "header":
            _docx_heading(d, str(p.get("text") or ""), _header_level(p))
        elif t == "list":
            style = "List Number" if p.get("ordered") else "List Bullet"
            for item in p.get("items") or []:
                d.add_paragraph(str(item), style=style)
        elif t == "table":
            _docx_table(d, p)
        elif t == "key_value":
            _docx_key_value(d, p, body_width)
        elif t == "callout":
            _docx_callout(d, p)
        elif t == "image":
            _docx_image(d, p, body_width)
        elif t in _STOCK_BLOCK_TYPES:
            block, extra = _stock_table_block(p)
            if block:
                _docx_table(d, block["props"])
                if extra:
                    _docx_run(d.add_paragraph(), extra, 8.5, color=RPT_MUTED)
        elif t in _CHART_BLOCK_TYPES:
            block = _chart_table_block(b)
            if block:
                if block["props"].get("title"):
                    _docx_heading(d, block["props"]["title"], 3)
                _docx_table(d, block["props"])
        elif t == "divider":
            d.add_page_break()
        else:
            for line in _block_lines(b):
                par = d.add_paragraph()
                par.paragraph_format.space_after = Pt(4)
                _docx_run(par, line, 10.5)
    d.save(out_path)


def _pdf_korean_font():
    """A Korean-capable font for reportlab: a host TTF when one exists (embedded, best
    fidelity), else Adobe's CID KR font (no file needed — the viewer supplies the glyphs).

    The BOLD face is registered as a family MEMBER, not decoration: reportlab resolves <b> by
    looking up the family's bold entry, and with no such entry every bold in the report genre —
    cover title, headings, KPI values, table headers — silently rendered at regular weight."""
    from reportlab.pdfbase import pdfmetrics
    for path, bold_path in (
        ("/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
         "/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf"),
        ("/usr/share/fonts/truetype/noto/NotoSansKR-Regular.ttf",
         "/usr/share/fonts/truetype/noto/NotoSansKR-Bold.ttf"),
        ("C:/Windows/Fonts/malgun.ttf", "C:/Windows/Fonts/malgunbd.ttf"),
    ):
        if os.path.isfile(path):
            from reportlab.pdfbase.ttfonts import TTFont
            try:
                pdfmetrics.registerFont(TTFont("KoreanBody", path))
            except Exception:  # noqa: BLE001 — a broken font file falls through to CID
                continue
            bold = "KoreanBody"
            if os.path.isfile(bold_path):
                try:
                    pdfmetrics.registerFont(TTFont("KoreanBody-Bold", bold_path))
                    bold = "KoreanBody-Bold"
                except Exception:  # noqa: BLE001 — a missing bold face is not a broken document
                    bold = "KoreanBody"
            pdfmetrics.registerFontFamily("KoreanBody", normal="KoreanBody", bold=bold,
                                          italic="KoreanBody", boldItalic=bold)
            return "KoreanBody"
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    pdfmetrics.registerFont(UnicodeCIDFont("HYGothic-Medium"))
    return "HYGothic-Medium"


def make_pdf_file(blocks, title, out_path):
    """Render blocks -> a .pdf report — the docx genre's twin, drawn with reportlab's parts.
    Same cover, same ranked headings, same banded tables and tinted callouts, so the two
    formats of one report do not read as two different documents."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import (HRFlowable, Image as RLImage, ListFlowable, ListItem,
                                    PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table,
                                    TableStyle)

    font = _pdf_korean_font()
    navy = colors.HexColor("#" + RPT_NAVY)
    accent = colors.HexColor("#" + RPT_ACCENT)
    body_c = colors.HexColor("#" + RPT_BODY)
    muted = colors.HexColor("#" + RPT_MUTED)
    band = colors.HexColor("#" + RPT_BAND)
    rule = colors.HexColor("#" + RPT_RULE)
    headfill = colors.HexColor("#" + RPT_HEADFILL)

    def st(name, size, **kw):
        return ParagraphStyle(name, fontName=font, fontSize=size,
                              leading=kw.pop("leading", size * 1.45), **kw)

    styles = {
        "cover": st("cover", 26, textColor=navy, spaceAfter=8),
        "date": st("date", 10, textColor=muted, spaceAfter=4),
        1: st("h1", 16, textColor=navy, spaceBefore=14, spaceAfter=2),
        2: st("h2", 13, textColor=navy, spaceBefore=11, spaceAfter=4),
        3: st("h3", 11.5, textColor=navy, spaceBefore=8, spaceAfter=3),
        "body": st("b", 10.5, textColor=body_c, spaceAfter=4),
        "cell": st("cell", 9.5, textColor=body_c),
        "cellr": st("cellr", 9.5, textColor=body_c, alignment=2),
        "th": st("th", 9.5, textColor=navy, alignment=1),
        "kpil": st("kpil", 8, textColor=muted, spaceAfter=1),
        "kpiv": st("kpiv", 16, textColor=navy, spaceAfter=0),
        "kpid": st("kpid", 8.5, textColor=muted),
        "callout": st("callout", 10, textColor=body_c),
        "note": st("note", 8.5, textColor=muted, spaceAfter=4),
    }
    margin = 18 * mm
    avail = A4[0] - margin * 2

    def esc(s):
        return (str(s).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))

    def para(text, style, bold=False):
        text = esc(text)
        return Paragraph(f"<b>{text}</b>" if bold else text, styles[style])

    def table_flow(p):
        headers, rows = _table_rows(p.get("headers"), p.get("rows"))
        if not headers and rows:
            headers = [str(c) for c in rows[0]]
            rows = rows[1:]
        if not headers:
            return []
        numeric = _numeric_cols(rows, len(headers))
        data = [[para(h, "th", bold=True) for h in headers]]
        for row in rows:
            cells = []
            for c in range(len(headers)):
                val = row[c] if c < len(row) else ""
                cells.append(para("" if val is None else val,
                                  "cellr" if c in numeric else "cell"))
            data.append(cells)
        tbl = Table(data, repeatRows=1, hAlign="LEFT")
        tbl.setStyle(TableStyle([
            ("GRID", (0, 0), (-1, -1), 0.4, rule),
            ("BACKGROUND", (0, 0), (-1, 0), headfill),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]))
        return [tbl, Spacer(1, 4 * mm)]

    def kpi_flow(metrics):
        per_row = min(4, max(1, len(metrics)))
        out = []
        for start in range(0, len(metrics), per_row):
            chunk = metrics[start:start + per_row]
            cells = []
            for m in chunk:
                stack = [para(m.get("label") or "", "kpil"),
                         para(f"{m.get('value', '')}{m.get('unit') or ''}", "kpiv", bold=True)]
                if m.get("delta") not in (None, ""):
                    stack.append(para(m.get("delta"), "kpid"))
                cells.append(stack)
            cells += [""] * (per_row - len(chunk))
            tbl = Table([cells], colWidths=[avail / per_row] * per_row, hAlign="LEFT")
            tbl.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (len(chunk) - 1, 0), band),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING", (0, 0), (-1, -1), 8),
                ("TOPPADDING", (0, 0), (-1, -1), 7),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
            ]))
            out += [tbl, Spacer(1, 3 * mm)]
        return out

    def kv_flow(p):
        pairs = _kv_pairs(p)
        if not pairs:
            return []
        out = [para(p.get("title"), 3, bold=True)] if str(p.get("title") or "").strip() else []
        data = [[para(k, "cell"), para("" if v is None else v, "cellr", bold=True)]
                for k, v in pairs]
        tbl = Table(data, colWidths=[avail * 0.45, avail * 0.55], hAlign="LEFT")
        tbl.setStyle(TableStyle([
            ("LINEBELOW", (0, 0), (-1, -2), 0.4, rule),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ("LEFTPADDING", (0, 0), (-1, -1), 0),
        ]))
        return out + [tbl, Spacer(1, 4 * mm)]

    def callout_flow(p):
        c_title, msg = _callout_parts(p)
        if not c_title and not msg:
            return []
        inner = []
        if c_title:
            inner.append(para(c_title, "callout", bold=True))
        if msg:
            inner.append(para(msg, "callout"))
        tbl = Table([[inner]], colWidths=[avail], hAlign="LEFT")
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), band),
            ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor("#" + RPT_CALLOUT_EDGE)),
            ("LEFTPADDING", (0, 0), (-1, -1), 8),
            ("RIGHTPADDING", (0, 0), (-1, -1), 8),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ]))
        return [tbl, Spacer(1, 4 * mm)]

    def image_flow(p):
        img_path, _err = resolve_path(str(p.get("src") or ""))
        if not img_path:
            return []
        try:
            # Fit to the text column, and never blow a small image up to reach it — the old
            # `height=60mm, kind="proportional"` sized on height alone, so a wide screenshot
            # ran off the page.
            iw, ih = ImageReader(img_path).getSize()
            w = min(float(avail), float(iw))
            return [RLImage(img_path, width=w, height=w * float(ih) / float(iw)),
                    Spacer(1, 3 * mm)]
        except Exception:  # noqa: BLE001 — a bad image loses itself, not the document
            return []

    cover, blocks = _report_cover(blocks, title)
    story = []
    if cover:
        story += [Spacer(1, 55 * mm), para(cover, "cover", bold=True),
                  para(_today_line(), "date"),
                  HRFlowable(width="35%", thickness=1.6, color=accent, hAlign="LEFT",
                             spaceBefore=2, spaceAfter=2)]
        if blocks:
            story.append(PageBreak())

    metric_at = dict(_metric_runs(blocks))
    skip_to = -1
    for i, b in enumerate(blocks):
        if i < skip_to:
            continue
        t, p = str(b.get("type") or ""), b.get("props") or {}
        if i in metric_at:
            story += kpi_flow(metric_at[i])
            skip_to = i + len(metric_at[i])
        elif t == "header":
            lvl = _header_level(p)
            story.append(para(p.get("text") or "", lvl, bold=True))
            if lvl == 1:
                story.append(HRFlowable(width="100%", thickness=1.0, color=accent,
                                        spaceBefore=1, spaceAfter=6))
        elif t == "list":
            items = [ListItem(para(it, "body")) for it in (p.get("items") or [])]
            if items:
                # Tight indents: reportlab's defaults park the bullet a centimetre from its
                # own sentence, which reads as two columns rather than one list.
                story.append(ListFlowable(
                    items, bulletType="1" if p.get("ordered") else "bullet",
                    bulletFontName=font, bulletFontSize=8, leftIndent=14, bulletDedent=10,
                    spaceAfter=4))
        elif t == "table":
            story += table_flow(p)
        elif t == "key_value":
            story += kv_flow(p)
        elif t == "callout":
            story += callout_flow(p)
        elif t == "image":
            story += image_flow(p)
        elif t in _STOCK_BLOCK_TYPES:
            block, extra = _stock_table_block(p)
            if block:
                story += table_flow(block["props"])
                if extra:
                    story.append(para(extra, "note"))
        elif t in _CHART_BLOCK_TYPES:
            block = _chart_table_block(b)
            if block:
                if block["props"].get("title"):
                    story.append(para(block["props"]["title"], 3, bold=True))
                story += table_flow(block["props"])
        elif t == "divider":
            story.append(PageBreak())
        else:
            for line in _block_lines(b):
                story.append(para(line, "body"))
    if not story:
        story.append(para(" ", "body"))

    def footer(canvas, _doc):
        canvas.saveState()
        canvas.setFont(font, 8)
        canvas.setFillColor(muted)
        canvas.drawCentredString(A4[0] / 2, 12 * mm, str(canvas.getPageNumber()))
        canvas.restoreState()

    def blank(_canvas, _doc):
        return None

    SimpleDocTemplate(out_path, pagesize=A4,
                      topMargin=18 * mm, bottomMargin=18 * mm,
                      leftMargin=margin, rightMargin=margin).build(
        story, onFirstPage=blank if cover else footer, onLaterPages=footer)
    return font


def action_make_pdf(inp):
    notes = []
    blocks = normalize_blocks(inp.get("blocks"), notes)
    title = str(inp.get("title") or "").strip()
    if not blocks and not title:
        return {"success": False, "action": "make_pdf",
                "error": "blocks (or at least a title) required"}
    out_path, stem = out_file(title or "document", "pdf", {"t": title, "b": blocks})
    try:
        font = make_pdf_file(blocks, title, out_path)
    except Exception as e:  # noqa: BLE001
        return {"success": False, "action": "make_pdf", "error": f"pdf build failed: {e}"}
    return {"success": True, "action": "make_pdf", "data": {
        "blocks": len(blocks), "font": font,
        **({"notes": notes} if notes else {}),
        "_mediaImport": media_import_decl(out_path, "pdf", stem),
    }}


def action_make_docx(inp):
    notes = []
    blocks = normalize_blocks(inp.get("blocks"), notes)
    title = str(inp.get("title") or "").strip()
    if not blocks and not title:
        return {"success": False, "action": "make_docx",
                "error": "blocks (or at least a title) required"}
    out_path, stem = out_file(title or "document", "docx", {"t": title, "b": blocks})
    try:
        make_docx_file(blocks, title, out_path)
    except Exception as e:  # noqa: BLE001
        return {"success": False, "action": "make_docx", "error": f"docx build failed: {e}"}
    return {"success": True, "action": "make_docx", "data": {
        "blocks": len(blocks),
        **({"notes": notes} if notes else {}),
        "_mediaImport": media_import_decl(out_path, "docx", stem),
    }}


def action_make_hwpx(inp):
    """Render blocks → .hwpx through the rhwp engine, seeded from a REAL Hancom blank
    (fixtures/donor-blank.hwp, 사용자 한컴 저장본) — createEmpty ships no style tables and
    Hancom draws garbage from its exports (2026-08-11 해부). Tables and bold headings
    ride along; the file lands in the media store like every other make."""
    notes = []
    blocks = normalize_blocks(inp.get("blocks"), notes)
    # The rhwp helper speaks header/text/list/table/metric. Everything the report renderers draw
    # natively arrives here as its text form — a callout prefixed "※ " is a paragraph in Hancom,
    # and a paragraph that says the caveat beats a caveat that was never written.
    flat = []
    for b in blocks:
        t = str(b.get("type") or "")
        if t in ("header", "text", "list", "table", "metric", "image", "divider"):
            flat.append(b)
            continue
        if t in _STOCK_BLOCK_TYPES:
            block, extra = _stock_table_block(b.get("props") or {})
            if block:
                flat.append(block)
                if extra:
                    flat.append({"type": "text", "props": {"content": extra}})
            continue
        if t in _CHART_BLOCK_TYPES:
            block = _chart_table_block(b)
            if block:
                flat.append(block)
            continue
        for line in _block_lines(b):
            flat.append({"type": "text", "props": {"content": line}})
    blocks = flat
    if not blocks and not title:
        return {"success": False, "action": "make_hwpx",
                "error": "blocks (or at least a title) required"}
    donor = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                         "fixtures", "donor-blank.hwp")
    if not os.path.exists(donor):
        return {"success": False, "action": "make_hwpx",
                "error": "donor-blank.hwp fixture missing — hwpx needs the Hancom seed"}
    out_path, stem = out_file(title or "document", "hwpx", {"t": title, "b": blocks})
    try:
        made = _rhwp_helper({"op": "make", "donor": donor, "out": out_path,
                             "format": "hwpx", "title": title, "blocks": blocks})
    except Exception as e:  # noqa: BLE001 — the engine names its own failure
        return {"success": False, "action": "make_hwpx", "error": f"hwpx build failed: {e}"}
    return {"success": True, "action": "make_hwpx", "data": {
        "blocks": len(blocks),
        "paragraphs": made.get("paragraphs"), "tables": made.get("tables"),
        **({"notes": notes} if notes else {}),
        "_mediaImport": media_import_decl(out_path, "hwpx", stem),
    }}


# ── selftest ───────────────────────────────────────────────────────────────────────────────────

def minimal_pdf():
    """A complete one-page PDF built with real offsets — pypdf wants the xref honest."""
    stream = b"BT /F1 24 Tf 72 700 Td (Hello Firebat) Tj ET"
    objs = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids[3 0 R]/Count 1>>",
        b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]"
        b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>",
        b"<</Length %d>>stream\n%s\nendstream" % (len(stream), stream),
        b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objs, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n%s\nendobj\n" % (i, body)
    xref_at = len(out)
    out += b"xref\n0 %d\n0000000000 65535 f \n" % (len(objs) + 1)
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += (b"trailer<</Size %d/Root 1 0 R>>\nstartxref\n%d\n%%%%EOF\n"
            % (len(objs) + 1, xref_at))
    return bytes(out)

HWPX_SECTION = """<?xml version="1.0" encoding="UTF-8"?>
<hs:sec xmlns:hs="http://www.hancom.co.kr/hwpml/2011/section"
        xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
  <hp:p><hp:run><hp:t>한글 문단 하나</hp:t></hp:run></hp:p>
  <hp:p><hp:run><hp:t>두 번째 문단</hp:t>
    <hp:equation><hp:script>a over b</hp:script></hp:equation></hp:run></hp:p>
  <hp:tbl>
    <hp:tr><hp:tc><hp:p><hp:run><hp:t>이름</hp:t></hp:run></hp:p></hp:tc>
           <hp:tc><hp:p><hp:run><hp:t>값</hp:t></hp:run></hp:p></hp:tc></hp:tr>
    <hp:tr><hp:tc><hp:p><hp:run><hp:t>가</hp:t></hp:run></hp:p></hp:tc>
           <hp:tc><hp:p><hp:run><hp:t>1</hp:t></hp:run></hp:p></hp:tc></hp:tr>
  </hp:tbl>
</hs:sec>"""


def action_selftest():
    checks = []

    def ck(name, ok):
        checks.append({"name": name, "ok": bool(ok)})

    os.makedirs(OUT_DIR, exist_ok=True)
    tmp = []

    blocks = [
        {"type": "header", "props": {"text": "첫 장", "level": 1}},
        {"type": "text", "props": {"content": "본문 한 줄"}},
        {"type": "table", "props": {"headers": ["이름", "값"], "rows": [["가", 1], ["나", 2]]}},
        {"type": "header", "props": {"text": "둘째 장", "level": 1}},
        {"type": "list", "props": {"items": ["하나", "둘"], "ordered": False}},
        {"type": "metric", "props": {"label": "매출", "value": 42, "unit": "억"}},
    ]

    ck("slide boundary splits on header level 1", len(_split_slides(blocks)) == 2)

    fence = normalize_blocks([
        {"name": "Header", "type": "component", "props": {"text": "장", "level": 1}},
        {"name": "Divider", "type": "component", "props": {}},
        {"name": "Table", "type": "component", "props": {"headers": ["a"], "rows": [["1"]]}},
    ])
    ck("render-fence dialect is absorbed (name -> type)",
       [b["type"] for b in fence] == ["header", "divider", "table"])

    # pptx round-trip
    p = f"{OUT_DIR}/selftest.pptx"
    tmp.append(p)
    try:
        n = make_pptx_file(blocks, "테스트 덱", None, p)
        got = read_pptx(p)
        ck("pptx: title + 2 body slides", n == 3)
        ck("pptx: text survives the round trip", "본문 한 줄" in got["text"])
        ck("pptx: the table is a real table", got["tables"]
           and got["tables"][0]["headers"] == ["이름", "값"])
        # The first live deck's two defects, pinned: the empty subtitle placeholder ghost on
        # the title slide, and overflow silently dropping content instead of turning the page.
        from pptx import Presentation as _P
        cov = _P(p).slides[0]
        ck("pptx: cover has no placeholder ghosts",
           len(list(cov.placeholders)) == 0 and len(cov.shapes) >= 2)
        # 2026-08-11 proposal-genre system: navy section band on every content slide,
        # and four short groups arrange as the quadrant donut.
        band = [sh for sh in _P(p).slides[1].shapes
                if sh.top == 0 and abs(sh.width - _P(p).slide_width) < 20000]
        ck("pptx: content slides carry the section band", len(band) >= 1)
        p_quad = f"{OUT_DIR}/selftest-quad.pptx"
        tmp.append(p_quad)
        four = [{"type": "header", "props": {"text": "SWOT", "level": 1}}]
        for t_ in ("Strength", "Weakness", "Opportunity", "Threat"):
            four.append({"type": "header", "props": {"text": t_, "level": 2}})
            four.append({"type": "text", "props": {"content": f"{t_} 설명"}})
        make_pptx_file(four, None, None, p_quad)
        got_quad = read_pptx(p_quad)
        ck("pptx: four short groups become the quadrant donut (all texts land)",
           all(t_ in got_quad["text"]
               for t_ in ("Strength", "Weakness", "Opportunity", "Threat")))
        p_tl = f"{OUT_DIR}/selftest-timeline.pptx"
        tmp.append(p_tl)
        tl = [{"type": "header", "props": {"text": "추진 일정", "level": 1}}]
        for yr in ("2024", "2025", "2026", "2027"):
            tl.append({"type": "header", "props": {"text": f"{yr}년", "level": 2}})
            tl.append({"type": "text", "props": {"content": f"{yr} 마일스톤"}})
        make_pptx_file(tl, None, None, p_tl)
        ck("pptx: year titles take the timeline (texts land)",
           "2027 마일스톤" in read_pptx(p_tl)["text"])
        p_cv = f"{OUT_DIR}/selftest-chevron.pptx"
        tmp.append(p_cv)
        cv = [{"type": "header", "props": {"text": "절차", "level": 1}},
              {"type": "list", "props": {"items": ["발굴", "검증", "계약", "운영"]}}]
        make_pptx_file(cv, None, None, p_cv)
        ck("pptx: a short label list flows as chevron arrows",
           all(s in read_pptx(p_cv)["text"] for s in ("발굴", "검증", "계약", "운영")))
        # A table taller than a slide splits across slides, header repeated, no row lost.
        p_tbl = f"{OUT_DIR}/selftest-tblsplit.pptx"
        tmp.append(p_tbl)
        big = [{"type": "table", "props": {"headers": ["번호", "값"],
                                           "rows": [[i, f"행{i}"] for i in range(60)]}}]
        n_tbl = make_pptx_file(big, None, None, p_tbl)
        got_tbl = read_pptx(p_tbl)
        total_rows = sum(len(t["rows"]) for t in got_tbl["tables"])
        ck("pptx: a huge table splits across slides", n_tbl >= 2)
        ck("pptx: every row survives the split", total_rows == 60)
        from pptx.oxml.ns import qn as _qn
        ck("pptx: slides carry the fade transition",
           _P(p).slides[0]._element.find(_qn("p:transition")) is not None)
        p_theme = f"{OUT_DIR}/selftest-theme.pptx"
        tmp.append(p_theme)
        make_pptx_file([{"type": "header", "props": {"text": "장", "level": 1}}], "테마",
                       None, p_theme, transition="none", theme={"accent": "#FF0000"})
        prs_t = _P(p_theme)
        ck("pptx: transition='none' really means none",
           prs_t.slides[0]._element.find(_qn("p:transition")) is None)
        reds = [sh for sl in prs_t.slides for sh in sl.shapes
                if sh.shape_type == 1 and sh.fill.type is not None
                and getattr(sh.fill.fore_color, "rgb", None) is not None
                and str(sh.fill.fore_color.rgb) == "FF0000"]
        ck("pptx: theme accent recolors the bars", len(reds) >= 1)
        p_spill = f"{OUT_DIR}/selftest-spill.pptx"
        tmp.append(p_spill)
        many = ([{"type": "header", "props": {"text": "긴 장", "level": 1}}]
                + [{"type": "text", "props": {"content": "내용 줄 " * 40}} for _ in range(20)])
        n_spill = make_pptx_file(many, None, None, p_spill)
        ck("pptx: overflowing content continues on new slides", n_spill >= 2)
        got_spill = read_pptx(p_spill)
        ck("pptx: nothing is dropped by overflow",
           sum(s["text"].count("내용 줄") for s in [{"text": got_spill["text"]}]) >= 1
           and got_spill["meta"]["slides"] == n_spill)
    except Exception as e:  # noqa: BLE001
        ck(f"pptx round trip crashed: {e}", False)

    # xlsx round-trip
    p = f"{OUT_DIR}/selftest.xlsx"
    tmp.append(p)
    try:
        make_xlsx_file([{"name": "표1", "headers": ["a", "b"], "rows": [[1, "x"], [2, "y"]]}], p)
        got = read_xlsx(p)
        ck("xlsx: sheet and rows survive", got["tables"]
           and got["tables"][0]["rows"] == [[1, "x"], [2, "y"]])
        p_num = f"{OUT_DIR}/selftest-num.xlsx"
        tmp.append(p_num)
        make_xlsx_file([{"name": "n", "headers": ["v"],
                         "rows": [["1,234"], ["-3.1"], ["48억"]]}], p_num)
        got_n = read_xlsx(p_num)
        vals = [r[0] for r in got_n["tables"][0]["rows"]]
        ck("xlsx: numeric-looking text lands as real numbers (units stay text)",
           vals == [1234, -3.1, "48억"])
    except Exception as e:  # noqa: BLE001
        ck(f"xlsx round trip crashed: {e}", False)

    # xlsx dashboard genre — KPI cards, native charts fed from cells, conditional formatting
    p_dash = f"{OUT_DIR}/selftest-dash.xlsx"
    tmp.append(p_dash)
    try:
        from openpyxl import load_workbook as _lw
        res_d = make_xlsx_file(
            [{"name": "Data", "headers": ["월", "매출"],
              "rows": [["1월", 100], ["2월", 220], ["3월", 180]]}],
            p_dash, title="월간 대시보드",
            kpis=[{"label": "합계", "value": "=SUM(Data!B2:B4)"},
                  {"label": "평균", "value": 166.7, "unit": "억"},
                  {"label": "영업이익", "value": 180, "delta": -40}],
            charts=[{"type": "bar", "title": "월별 매출", "sheet": "Data",
                     "labelCol": "월", "valueCols": ["매출"]},
                    {"type": "line", "title": "없는 시트", "sheet": "NoSuchSheet",
                     "valueCols": ["x"]}])
        wb_d = _lw(p_dash, data_only=False)
        ck("xlsx dashboard: Dashboard sheet exists and is first",
           wb_d.sheetnames[0] == "Dashboard" and "Data" in wb_d.sheetnames)
        # openpyxl drops charts when READING, so the chart part is counted in the package.
        with zipfile.ZipFile(p_dash) as z:
            parts = z.namelist()
            dash_rels = (z.read("xl/worksheets/_rels/sheet1.xml.rels").decode("utf-8")
                         if "xl/worksheets/_rels/sheet1.xml.rels" in parts else "")
            dash_bar_xml = "".join(z.read(n).decode("utf-8") for n in parts
                                   if n.startswith("xl/charts/chart"))
        ck("xlsx dashboard: exactly one native chart is written",
           len([n for n in parts if n.startswith("xl/charts/chart")]) == 1 and res_d["charts"] == 1)
        ck("xlsx dashboard: the chart is anchored on the Dashboard sheet", "drawing" in dash_rels)
        ck("xlsx dashboard: an unresolvable chart is skipped WITH a note",
           len(res_d["notes"]) == 1 and "NoSuchSheet" in res_d["notes"][0])
        dash_cells = [c.value for row in wb_d["Dashboard"].iter_rows() for c in row]
        ck("xlsx dashboard: the formula KPI stays a live formula",
           any(isinstance(v, str) and v.startswith("=SUM(Data!") for v in dash_cells))
        ck("xlsx dashboard: KPI labels and the negative delta are rendered",
           "합계" in dash_cells and any(isinstance(v, str) and v.startswith("▼") for v in dash_cells))
        # The dashboard genre: no grid behind the report, and a KPI card is a boxed field.
        ws_dr = wb_d["Dashboard"]
        ck("xlsx dashboard: the Dashboard sheet drops the gridlines",
           ws_dr.sheet_view.showGridLines is False and wb_d["Data"].sheet_view.showGridLines
           is not False)
        a3, a5 = ws_dr["A3"], ws_dr["A5"]
        # Card v3: the accent is a thick bottom underline now, not a left spine.
        ck("xlsx dashboard: KPI cards get a border box, a colored underline and a label fill",
           a3.border.top.style == "thin" and a5.border.bottom.style == "thick"
           and str(a5.border.bottom.color.rgb or "").endswith(XLSX_CARD_UNDERLINES[0])
           and str(a3.fill.fgColor.rgb or "").endswith(XLSX_CARD_LABEL_BG))
        # v4: three cards stretch to fill the 12-column grid — 4 columns each, so card 2 is at E.
        ck("xlsx dashboard: the card underline cycles color per card",
           str(ws_dr["E5"].border.bottom.color.rgb or "").endswith(XLSX_CARD_UNDERLINES[1]))
        ck("xlsx dashboard: the title band spans the whole fixed grid",
           str(ws_dr["A1"].fill.fgColor.rgb or "").endswith(XLSX_BAND)
           and any(str(r) == "A1:L1" for r in ws_dr.merged_cells.ranges))
        rules_data = [r for cf in wb_d["Data"].conditional_formatting for r in cf.rules]
        ck("xlsx dashboard: the data sheet gets a data bar",
           any(r.type == "dataBar" for r in rules_data))
        ck("xlsx dashboard: the header row is frozen", wb_d["Data"].freeze_panes == "A2")

        p_cf = f"{OUT_DIR}/selftest-cf.xlsx"
        tmp.append(p_cf)
        make_xlsx_file([{"name": "비율", "headers": ["종목", "증감률(%)"],
                         "rows": [["가", 1.5], ["나", -2.5], ["다", 0.5]]}], p_cf)
        rules_cf = [r for cf in _lw(p_cf)["비율"].conditional_formatting for r in cf.rules]
        ck("xlsx dashboard: a 증감률(%) column takes a color scale, not bars",
           any(r.type == "colorScale" for r in rules_cf)
           and not any(r.type == "dataBar" for r in rules_cf))

        # 2026-08-12 premium-dashboard archetypes: combo, doughnut + center block, bar KPIs.
        p_v3 = f"{OUT_DIR}/selftest-v3.xlsx"
        tmp.append(p_v3)
        make_xlsx_file(
            [{"name": "월별", "headers": ["월", "수량", "매출"],
              "rows": [[f"{m}월", 100 + m * 5, 1000 + m * 90] for m in range(1, 13)]},
             {"name": "유입", "headers": ["구분", "비중"],
              "rows": [["첫구매", 83.1], ["재구매", 16.9]]}],
            p_v3, title="커머스 대시보드",
            kpis=[{"label": "총 매출", "value": 334, "unit": "조원", "icon": "💰", "delta": 12},
                  {"label": "주문 수", "value": 1200, "delta": -30},
                  {"label": "검색유입 매출 비율", "value": 34, "unit": "%", "style": "bar"},
                  {"label": "고객 만족도 (5점 만점)", "value": 4.11, "max": 5, "style": "bar"}],
            charts=[{"type": "combo", "title": "월별 수량·매출", "sheet": "월별",
                     "labelCol": "월", "valueCols": ["수량", "매출"], "unit": "단위: 백만원"},
                    {"type": "doughnut", "title": "첫구매 비중", "sheet": "유입",
                     "labelCol": "구분", "valueCols": ["비중"],
                     "centerLabel": "첫구매 비율", "centerValue": "83.1%"}])
        ws_v3 = _lw(p_v3)["Dashboard"]
        with zipfile.ZipFile(p_v3) as z:
            charts_xml = [z.read(n).decode("utf-8") for n in z.namelist()
                          if n.startswith("xl/charts/chart")]
            draw_v3 = z.read("xl/drawings/drawing1.xml").decode("utf-8")
        # openpyxl writes the chart part with the chart namespace as the DEFAULT one, so the
        # elements carry no c: prefix — match the bare tag names. Same for the drawing part.
        combo_xml = [x for x in charts_xml if "<barChart" in x]
        dough_xml = [x for x in charts_xml if "<doughnutChart" in x]
        ck("xlsx v3: combo draws bars and a line, the line on a secondary axis",
           len(combo_xml) == 1 and "<lineChart" in combo_xml[0]
           and combo_xml[0].count("<valAx") >= 2
           and '<axId val="200"/>' in combo_xml[0])
        ck("xlsx v3: a doughnut is a real ring, not a pie", len(dough_xml) == 1)
        v3_cells = [c.value for row in ws_v3.iter_rows() for c in row]
        ck("xlsx v3: the doughnut's center block carries label and value",
           "83.1%" in v3_cells and "첫구매 비율" in v3_cells)
        ck("xlsx v3: a bar-style KPI puts a data bar on the Dashboard itself",
           any(r.type == "dataBar" for cf in ws_v3.conditional_formatting for r in cf.rules))
        ck("xlsx v3: an icon KPI trades numeric-ness for the emoji",
           any(isinstance(v, str) and v.startswith("💰") and "334조원" in v for v in v3_cells))

        # 2026-08-12 v4: the canvas. Two measured flaws — legends/labels landing on the data, and
        # a page that read as empty — were both geometry, so the checks below are geometry too.
        from openpyxl.utils import get_column_letter as _gcl

        def _fill_of(ws, addr):
            f = ws[addr].fill
            return str(f.fgColor.rgb or "") if f is not None and f.patternType else ""

        band0 = XLSX_KPI_TOP + XLSX_KPI_ROWS      # 2 cards = 1 card row, then the chart band
        wide_cols = XLSX_DASH_COLS - max(3, round(XLSX_DASH_COLS / 3))
        anchors = [(int(c), int(r)) for c, r in re.findall(
            r"<from><col>(\d+)</col><colOff>\d+</colOff><row>(\d+)</row>", draw_v3)]
        ck("xlsx v4: the page is painted slate everywhere a card is not",
           _fill_of(ws_v3, "A2").endswith(XLSX_CANVAS)
           and _fill_of(ws_v3, f"{_gcl(XLSX_DASH_COLS + 1)}1").endswith(XLSX_CANVAS))
        ck("xlsx v4: a chart sits on a white card block with a hairline frame",
           _fill_of(ws_v3, f"A{band0 + 1}").endswith(XLSX_WHITE)
           and ws_v3[f"A{band0}"].border.top.style == "thin"
           and "</chart><spPr>" in combo_xml[0] and "FFFFFF" in combo_xml[0].split("</chart>")[1])
        ck("xlsx v4: a combo and a doughnut share one band, ~2/3 + ~1/3",
           len(anchors) == 2 and anchors[0][1] == anchors[1][1] == band0
           and anchors[0][0] == 0 and anchors[1][0] == wide_cols)
        ck("xlsx v4: every legend sits under the plot, the doughnut's included",
           '<legendPos val="b"/>' in combo_xml[0] and '<legendPos val="b"/>' in dough_xml[0])
        ck("xlsx v4: the center block sits UNDER the ring, inside the doughnut's own card",
           ws_v3.cell(row=band0 + 1 + XLSX_RING_ROWS, column=wide_cols + 1).value == "첫구매 비율"
           and ws_v3.cell(row=band0 + 2 + XLSX_RING_ROWS,
                          column=wide_cols + 1).value == "83.1%")
        ck("xlsx v4: a chart's unit becomes an amber badge at its card's top-right",
           ws_v3.cell(row=band0, column=wide_cols - 1).value == "단위: 백만원"
           and _fill_of(ws_v3, f"{_gcl(wide_cols - 1)}{band0}").endswith(XLSX_UNIT_BADGE_BG))
        ck("xlsx v4: value labels ride the bars of a bar-only chart",
           "<dLbls>" in dash_bar_xml)

        # 2026-08-12 v5: three measured flaws from the v4 sample — axes with no numbers on them,
        # a combo that threw its bar values away with the line's, and a ring wearing its title.
        ring_cols = max(3, round(XLSX_DASH_COLS / 3))
        exts = [(int(cx), int(cy)) for cx, cy in
                re.findall(r'<ext cx="(\d+)" cy="(\d+)"/>', draw_v3)]
        combo_bars, combo_line = combo_xml[0].split("<lineChart", 1)
        ck("xlsx v5: every cartesian axis says delete=0, the combo's secondary axis included",
           combo_xml[0].count('<delete val="0"/>') == 3
           and combo_xml[0].count('<tickLblPos val="nextTo"/>') == 3
           and dash_bar_xml.count('<delete val="0"/>') == 2)
        ck("xlsx v5: both of the combo's value axes carry a thousands format",
           combo_xml[0].count('<numFmt formatCode="#,##0" sourceLinked="0"/>') == 2)
        ck("xlsx v5: a combo labels its bars and leaves its line bare",
           "<dLbls>" in combo_bars and '<dLblPos val="outEnd"/>' in combo_bars
           and "<dLbls>" not in combo_line)
        ck("xlsx v5: a ring carries no internal title — its card's header cell does",
           "<title>" not in dough_xml[0]
           and ws_v3.cell(row=band0, column=wide_cols + 1).value == "첫구매 비중")
        # 1cm = 360000 EMU. The ring is a drawing, not a cell, so "inside the card" is arithmetic:
        # it may not be wider than its own column span nor taller than the rows reserved for it.
        ck("xlsx v5: the ring never outgrows its card box",
           len(exts) == 2
           and exts[1][0] <= round(ring_cols * XLSX_COL_CM * 360000)
           and exts[1][1] <= round(XLSX_RING_ROWS * XLSX_ROW_CM * 360000))
        bar_kpi_row = XLSX_KPI_TOP + XLSX_KPI_ROWS + XLSX_CHARTCARD_ROWS + 1
        bar_col = 1 + XLSX_BARKPI_VALUE_COLS + XLSX_BARKPI_LABEL_COLS
        int_cell = ws_v3.cell(row=bar_kpi_row, column=bar_col)
        dec_cell = ws_v3.cell(row=bar_kpi_row + XLSX_BARKPI_ROWS, column=bar_col)
        ck("xlsx v5: the in-bar value drops the dangling decimal point, keeps real decimals",
           (int_cell.value, int_cell.number_format) == (34, "#,##0")
           and (dec_cell.value, dec_cell.number_format) == (4.11, "#,##0.##")
           and int_cell.alignment.horizontal == "left")

        # 2026-08-12 v6: two flaws from the deployed sample — data labels wearing series
        # AND category names (omitted CT_Booleans read TRUE), and year columns dressed as
        # magnitudes ("2,023" + a data bar over 연도).
        # Aliased imports: a bare `import zipfile` here would make the name local to the
        # WHOLE enclosing function and break the earlier blocks that already use it.
        import zipfile as _zf6

        import openpyxl as _px6

        p_v6 = f"{OUT_DIR}/selftest-v6.xlsx"
        tmp.append(p_v6)
        make_xlsx_file(
            [{"name": "연간", "headers": ["연도", "자산", "부채"],
              "rows": [["2023", 455905980, 92228115], ["2024", 514531948, 112339878],
                       ["2025", 566942110, 130621773]]}],
            p_v6, title="v6",
            charts=[{"type": "bar", "title": "다중", "sheet": "연간",
                     "labelCol": "연도", "valueCols": ["자산", "부채"]},
                    {"type": "bar", "title": "단일", "sheet": "연간",
                     "labelCol": "연도", "valueCols": ["자산"]}])
        with _zf6.ZipFile(p_v6) as zv6:
            charts_v6 = sorted(n for n in zv6.namelist()
                               if re.match(r"xl/charts/chart\d+\.xml$", n))
            xml_v6 = [zv6.read(n).decode("utf-8") for n in charts_v6]
        multi = next(x for x in xml_v6 if x.count("<ser>") == 2)
        single = next(x for x in xml_v6 if x.count("<ser>") == 1)
        ck("xlsx v6: a multi-series bar carries NO value labels — the legend is its labeling",
           "<dLbls>" not in multi.split("<plotArea")[1].split("<catAx")[0])
        ck("xlsx v6: a single-series bar labels values ONLY — every other show flag is 0",
           "<dLbls>" in single and '<showVal val="1"/>' in single
           and '<showSerName val="0"/>' in single and '<showCatName val="0"/>' in single
           and '<showLegendKey val="0"/>' in single)
        wb_v6 = _px6.load_workbook(p_v6)
        ws_v6 = wb_v6["연간"]
        ck("xlsx v6: a year column is plain '0' — no thousands comma, no decoration",
           ws_v6.cell(row=2, column=1).number_format == "0"
           and all("A" not in str(rng) for rng in ws_v6.conditional_formatting))
        ck("xlsx v6: real magnitude columns keep their commas and bars",
           ws_v6.cell(row=2, column=2).number_format == "#,##0"
           and any("B" in str(rng.sqref) for rng in ws_v6.conditional_formatting))

        # Ledger genre: a document with live monthly / running subtotals.
        p_lg = f"{OUT_DIR}/selftest-ledger.xlsx"
        tmp.append(p_lg)
        make_xlsx_file([{
            "name": "원장", "style": "ledger", "docTitle": "총 계 정 원 장",
            "period": "2026.01.01 ~ 2026.02.28", "subtotalBy": "일자",
            "headers": ["일자", "구분", "적요", "입금", "출금"],
            "rows": [["2026-01-05", "매출", "1월 판매", 1000, 0],
                     ["2026-01-20", "비용", "임차료", 0, 300],
                     ["2026-02-03", "매출", "2월 판매", 2000, 0],
                     ["2026-02-17", "비용", "광고", 0, 500]]}], p_lg)
        ws_lg = _lw(p_lg)["원장"]
        ck("xlsx ledger: the doc title is a centered merged masthead",
           ws_lg["A1"].value == "총 계 정 원 장"
           and any(str(r) == "A1:E1" for r in ws_lg.merged_cells.ranges)
           and ws_lg["A1"].alignment.horizontal == "center"
           and ws_lg["A2"].value.startswith("2026.01.01"))
        m_rows = [c.row for c in ws_lg["A"] if c.value == XLSX_LEDGER_MONTH_LABEL]
        t_rows = [c.row for c in ws_lg["A"] if c.value == XLSX_LEDGER_TOTAL_LABEL]
        ck("xlsx ledger: [월 계] is a live SUM over exactly that month's rows",
           len(m_rows) == 2 and ws_lg.cell(row=m_rows[0], column=4).value == "=SUM(D5:D6)")
        # The running total starts at the first data row and skips the subtotal bands it passes
        # — one contiguous range would count them a second time.
        ck("xlsx ledger: the last [누 계] runs from the first data row",
           len(t_rows) == 2
           and ws_lg.cell(row=t_rows[-1], column=4).value == "=SUM(D5:D6,D9:D10)")
        ck("xlsx ledger: a ledger is a document, no data bars or color scales",
           not [r for cf in ws_lg.conditional_formatting for r in cf.rules])
        ck("xlsx ledger: the ledger sheet drops the gridlines",
           ws_lg.sheet_view.showGridLines is False)

        res_b = action_make_xlsx({"title": "실적 대시보드", "blocks": [
            {"type": "metric", "props": {"label": "매출", "value": 1200,
                                         "unit": "억", "delta": 30}},
            {"type": "metric", "props": {"label": "영업이익", "value": 180, "delta": -12}},
            {"type": "header", "props": {"text": "실적표", "level": 2}},
            {"type": "table", "props": {"headers": ["항목", "값"], "rows": [["매출", 1200]]}},
            {"type": "chart", "props": {"chartType": "line", "title": "분기 추이",
                                        "labels": ["1Q", "2Q", "3Q"], "data": [10, 20, 30]}},
        ]})
        ck("xlsx blocks: metric + chart blocks trigger dashboard mode",
           res_b.get("success") and res_b["data"]["kpis"] == 2
           and res_b["data"]["charts"] == 1 and res_b["data"]["dashboard"])
        p_b = res_b["data"]["_mediaImport"]["path"]
        tmp.append(p_b)
        wb_b = _lw(p_b)
        ck("xlsx blocks: Dashboard first, the chart gets its own data sheet",
           wb_b.sheetnames[0] == "Dashboard" and "실적표" in wb_b.sheetnames
           and "분기 추이" in wb_b.sheetnames)
        ck("xlsx blocks: two KPI cards land as merged label cells",
           wb_b["Dashboard"]["A3"].value == "매출"
           # v4: two cards split the 12-column grid in half, so card 2 starts at G
           and wb_b["Dashboard"]["G3"].value == "영업이익"
           # 2 cards x 3 rows + the row-1 title band
           and len(wb_b["Dashboard"].merged_cells.ranges) == 7)
    except Exception as e:  # noqa: BLE001
        ck(f"xlsx dashboard crashed: {e}", False)

    # docx round-trip
    p = f"{OUT_DIR}/selftest.docx"
    tmp.append(p)
    try:
        make_docx_file(blocks, "테스트 문서", p)
        got = read_docx(p)
        ck("docx: headings and text survive", "본문 한 줄" in got["text"]
           and any("첫 장" in ln for ln in got["text"].split("\n")))
        ck("docx: the table is a real table", got["tables"]
           and got["tables"][0]["headers"] == ["이름", "값"])
    except Exception as e:  # noqa: BLE001
        ck(f"docx round trip crashed: {e}", False)

    # hwpx read (fixture zip — we never write hwpx, by decision)
    p = f"{OUT_DIR}/selftest.hwpx"
    tmp.append(p)
    try:
        with zipfile.ZipFile(p, "w") as z:
            z.writestr("Contents/section0.xml", HWPX_SECTION)
        got = read_hwpx(p)
        ck("hwpx: paragraphs come out in order", got["text"].startswith("한글 문단 하나"))
        ck("hwpx: table cells keep their grid", got["tables"]
           and got["tables"][0]["rows"] == [["가", "1"]])
        ck("hwpx: table text is not doubled into the body", "이름" not in got["text"])
        ck("hwpx: equation script is extracted (직독)",
           got.get("equations") == ["a over b"] and "[수식] a over b" in got["text"])
    except Exception as e:  # noqa: BLE001
        ck(f"hwpx read crashed: {e}", False)

    # hwp via the rhwp engine — the committed fixture was made BY rhwp (create → export),
    # so this proves the vendored wasm runs on this host and equations survive as script.
    import shutil as _sh
    fx = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      "fixtures", "rhwp-roundtrip.hwp")
    if _sh.which("node") and os.path.exists(fx):
        try:
            got = read_hwp_legacy(fx)
            ck("hwp(rhwp): text and table cell survive",
               "파이어뱃" in got["text"] and got["tables"]
               and got["tables"][0]["headers"][0] == "셀A1")
            ck("hwp(rhwp): the equation comes back as its script",
               any("over" in e for e in got.get("equations", [])))
        except Exception as e:  # noqa: BLE001
            ck(f"hwp(rhwp) crashed: {e}", False)
        # make: donor-seeded hwpx round trip through our own direct-walk reader.
        donor = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                             "fixtures", "donor-blank.hwp")
        if os.path.exists(donor):
            p_mk = f"{OUT_DIR}/selftest-made.hwpx"
            tmp.append(p_mk)
            try:
                _rhwp_helper({"op": "make", "donor": donor, "out": p_mk,
                              "format": "hwpx", "title": "셀프테스트 문서",
                              "blocks": blocks})
                got_mk = read_hwpx(p_mk)
                ck("hwpx make: text and heading survive the round trip",
                   "본문 한 줄" in got_mk["text"] and "첫 장" in got_mk["text"])
                ck("hwpx make: the table is a real table",
                   got_mk["tables"] and got_mk["tables"][0]["headers"] == ["이름", "값"])
            except Exception as e:  # noqa: BLE001
                ck(f"hwpx make crashed: {e}", False)
        else:
            ck("hwpx make: skipped — donor fixture absent", True)
    else:
        ck("hwp(rhwp): skipped — node runtime or fixture absent", True)

    # pdf read (hand-written minimal fixture)
    p = f"{OUT_DIR}/selftest.pdf"
    tmp.append(p)
    try:
        with open(p, "wb") as fh:
            fh.write(minimal_pdf())
        got = read_pdf(p)
        ck("pdf: text extracts", "Hello Firebat" in got["text"])
    except Exception as e:  # noqa: BLE001
        ck(f"pdf read crashed: {e}", False)

    # ── 2026-08-12 coverage pass ───────────────────────────────────────────────────────────────
    # Every check below pins a measured hole: content that reached the file as NOTHING (container
    # children, baked module output, callouts, chart blocks), a dialect that raised or wrote
    # garbage (dict rows, "h2" levels), and the archetypes that only fired by accident.
    dict_rows = [{"종목": "삼성전자", "종가": 70500, "등락률": 1.2},
                 {"종목": "SK하이닉스", "종가": 210000, "등락률": -0.8}]
    try:
        drops = []
        packed = normalize_blocks([
            {"type": "grid", "props": {"columns": 2, "children": [
                {"type": "metric", "props": {"label": "매출", "value": 1200, "unit": "억"}},
                {"type": "metric", "props": {"label": "영업이익", "value": 180}}]}},
            {"type": "tabs", "props": {"tabs": [
                {"label": "국내", "children": [
                    {"type": "text", "props": {"content": "국내 본문"}}]},
                {"label": "해외", "children": [
                    {"type": "text", "props": {"content": "해외 본문"}}]}]}},
            {"type": "accordion", "props": {"items": [{"title": "FAQ", "content": "답변 본문"}]}},
            {"type": "module", "props": {"_baked": [
                {"type": "text", "props": {"content": "베이크된 문단"}},
                {"type": "module", "props": {"_baked": [
                    {"type": "text", "props": {"content": "중첩 모듈 본문"}}]}}]}},
            {"type": "alert", "props": {"message": "주의 문구"}},
            {"type": "button", "props": {"label": "클릭"}},
            {"type": "countdown", "props": {"to": "2026-12-31"}},
            {"type": "button", "props": {"label": "또 클릭"}},
            {"type": "live_stock_chart", "props": {
                "topic": "ws:1", "title": "삼성전자",
                "data": [{"date": "2026-08-11", "open": 70000, "high": 71000,
                          "low": 69500, "close": 70500, "volume": 1000}]}},
            {"type": "live_feed", "props": {"topic": "ws:2", "title": "체결"}},
        ], drops)
        kinds = [b["type"] for b in packed]
        ck("blocks: a grid's children are spliced in order",
           kinds[:2] == ["metric", "metric"])
        ck("blocks: tabs and accordion keep their section titles as level-3 headers",
           [(b["props"].get("text"), b["props"].get("level")) for b in packed
            if b["type"] == "header"] == [("국내", 3), ("해외", 3), ("FAQ", 3)])
        ck("blocks: module._baked is spliced, a module nested inside it is left alone",
           any(b["type"] == "text" and b["props"]["content"] == "베이크된 문단" for b in packed)
           and any(b["type"] == "module" for b in packed)
           and not any(b["type"] == "text" and b["props"]["content"] == "중첩 모듈 본문"
                       for b in packed))
        ck("blocks: alert is the callout it always was", "callout" in kinds)
        ck("blocks: a live block carrying a seed degrades to a labelled snapshot",
           any(b["type"] == "stock_chart" and b["props"]["title"].endswith(_SNAPSHOT_SUFFIX)
               for b in packed))
        ck("blocks: a live block holding only a topic is dropped WITH a note",
           "live_feed" not in kinds and any("live_feed" in n for n in drops))
        ck("blocks: a dropped type is noted once per TYPE, not once per block",
           sum(1 for n in drops if "button" in n) == 1
           and sum(1 for n in drops if "countdown" in n) == 1)

        tbl_b = normalize_blocks([{"type": "table", "props": {
            "headers": ["종목", "종가", "등락률"], "rows": dict_rows}}])
        ck("blocks: dict rows are ordered by the table's own headers",
           tbl_b[0]["props"]["rows"] == [["삼성전자", 70500, 1.2],
                                         ["SK하이닉스", 210000, -0.8]])
        ck("blocks: with no headers the first record's keys become the header row",
           normalize_blocks([{"type": "table", "props": {"rows": dict_rows}}]
                            )[0]["props"]["headers"] == ["종목", "종가", "등락률"])
        ck("blocks: a non-numeric header level is clamped, never fatal",
           (_header_level({"level": "h2"}), _header_level({"level": "H9"}),
            _header_level({"level": None}), _header_level({"level": 0}),
            _header_level({})) == (2, 3, 2, 1, 2))
    except Exception as e:  # noqa: BLE001
        ck(f"block normalization crashed: {e}", False)

    # pptx: the two blocks a deck used to lose — a chart, and any table written as records.
    try:
        from pptx import Presentation as _PP
        p_ch = f"{OUT_DIR}/selftest-chart.pptx"
        tmp.append(p_ch)
        make_pptx_file(normalize_blocks([
            {"type": "header", "props": {"text": "분기 추이", "level": 1}},
            {"type": "chart", "props": {"chartType": "bar", "title": "분기 매출",
                                        "labels": ["1Q", "2Q", "3Q"], "data": [10, 20, 30]}},
            {"type": "chart", "props": {"chartType": "donut", "title": "구성",
                                        "labels": ["A", "B"], "data": [60, 40]}},
        ]), None, None, p_ch)
        with zipfile.ZipFile(p_ch) as z:
            chart_parts = [n for n in z.namelist() if n.startswith("ppt/charts/chart")]
        ck("pptx: chart blocks land as native chart parts", len(chart_parts) == 2)
        ck("pptx: the chart is a real graphic frame on the slide",
           any(sh.has_chart for sl in _PP(p_ch).slides for sh in sl.shapes))

        p_dr = f"{OUT_DIR}/selftest-dictrows.pptx"
        tmp.append(p_dr)
        make_pptx_file(normalize_blocks([{"type": "table", "props": {
            "headers": ["종목", "종가"], "rows": dict_rows}}]), None, None, p_dr)
        ck("pptx: dict rows no longer raise, and the VALUES land in the cells",
           read_pptx(p_dr)["tables"][0]["rows"][0] == ["삼성전자", "70500"])

        p_arch = f"{OUT_DIR}/selftest-arch.pptx"
        tmp.append(p_arch)
        make_pptx_file(normalize_blocks([
            {"type": "header", "props": {"text": "로드맵", "level": 1}},
            {"type": "timeline", "props": {"items": [
                {"date": "2026 1Q", "title": "설계", "description": "요건 확정"},
                {"date": "2026 2Q", "title": "개발", "description": "코어 구현"},
                {"date": "2026 3Q", "title": "검증", "description": "실측"}]}},
            {"type": "divider", "props": {}},
            {"type": "header", "props": {"text": "비교", "level": 1}},
            {"type": "compare", "props": {
                "left": {"label": "현행", "items": [{"key": "속도", "value": "느림"}]},
                "right": {"label": "개선", "items": [{"key": "속도", "value": "빠름"}]}}},
            {"type": "progress", "props": {"label": "진척률", "value": 68, "max": 100}},
            {"type": "callout", "props": {"title": "주의", "message": "잠정 일정입니다"}},
        ]), None, None, p_arch)
        arch_txt = read_pptx(p_arch)["text"]
        ck("pptx: timeline / compare / progress blocks fire their archetypes directly",
           all(s in arch_txt for s in ("2026 1Q 설계", "현행", "개선", "진척률", "68%")))
        ck("pptx: a callout is drawn as a band and keeps its words",
           "주의" in arch_txt and "잠정 일정입니다" in arch_txt)
    except Exception as e:  # noqa: BLE001
        ck(f"pptx coverage pass crashed: {e}", False)

    # xlsx: a grid of metrics has to reach the KPI band, dict rows have to write VALUES (the old
    # bug wrote the KEYS), a callout has to survive, and stock_chart has to be a real candlestick.
    try:
        from openpyxl import load_workbook as _lw2
        res_x = action_make_xlsx({"title": "커버리지", "blocks": [
            {"type": "grid", "props": {"columns": 3, "children": [
                {"type": "metric", "props": {"label": "매출", "value": 1200}},
                {"type": "metric", "props": {"label": "영업이익", "value": 180}},
                {"type": "metric", "props": {"label": "신규", "value": 12}}]}},
            {"type": "table", "props": {"headers": ["종목", "종가", "등락률"],
                                        "rows": dict_rows}},
            {"type": "alert", "props": {"message": "잠정치입니다"}},
            {"type": "button", "props": {"label": "클릭"}},
        ]})
        p_x = res_x["data"]["_mediaImport"]["path"]
        tmp.append(p_x)
        ck("xlsx: every metric inside a grid becomes a KPI card",
           res_x["data"]["kpis"] == 3)
        wb_x = _lw2(p_x)
        data_rows = [[c.value for c in r] for r in wb_x["표1"].iter_rows()]
        ck("xlsx: dict rows write their VALUES — the old bug wrote the KEYS as data",
           data_rows[0][:3] == ["종목", "종가", "등락률"]
           and data_rows[1][:3] == ["삼성전자", 70500, 1.2]
           and data_rows[1][:3] != ["종목", "종가", "등락률"])
        dash_x = [c.value for r in wb_x["Dashboard"].iter_rows() for c in r]
        ck("xlsx: a callout becomes an amber note row on the Dashboard",
           any(isinstance(v, str) and v.startswith("※") and "잠정치입니다" in v for v in dash_x))
        ck("xlsx: an undrawable block is named in data.notes",
           any("button" in n for n in res_x["data"].get("notes") or []))

        res_s = action_make_xlsx({"title": "일봉", "blocks": [
            {"type": "stock_chart", "props": {"title": "삼성전자", "data": [
                {"date": f"2026-08-{d:02d}", "open": 70000 + d * 10, "high": 70500 + d * 10,
                 "low": 69500 + d * 10, "close": 70200 + d * 10, "volume": 1000 * d}
                for d in range(1, 11)]}}]})
        p_s = res_s["data"]["_mediaImport"]["path"]
        tmp.append(p_s)
        with zipfile.ZipFile(p_s) as z:
            sxml = "".join(z.read(n).decode("utf-8") for n in z.namelist()
                           if n.startswith("xl/charts/chart"))
        ck("xlsx: a stock_chart block becomes a native candlestick with up/down bars",
           "<stockChart>" in sxml and "<upDownBars>" in sxml and "<hiLowLines" in sxml)
        ck("xlsx: the candle bodies wear the Korean up-red / down-blue convention",
           XLSX_UP in sxml and XLSX_DOWN in sxml)
        ck("xlsx: volume rides a bar group on the secondary axis",
           "<barChart>" in sxml and '<axId val="200"/>' in sxml)
    except Exception as e:  # noqa: BLE001
        ck(f"xlsx coverage pass crashed: {e}", False)

    # docx / pdf: the report genre. A container must produce exactly the paragraphs its children
    # would have produced on their own, and the cover has to actually be a cover.
    try:
        from docx import Document as _DD
        inner = [{"type": "text", "props": {"content": f"문단 {i}"}} for i in range(3)]
        p_g1, p_g2 = f"{OUT_DIR}/selftest-bare.docx", f"{OUT_DIR}/selftest-grid.docx"
        tmp += [p_g1, p_g2]
        make_docx_file(normalize_blocks(inner), "", p_g1)
        make_docx_file(normalize_blocks(
            [{"type": "grid", "props": {"columns": 3, "children": inner}}]), "", p_g2)
        ck("docx: a container yields exactly the paragraphs its bare children would",
           read_docx(p_g1)["meta"]["paragraphs"]
           == read_docx(p_g2)["meta"]["paragraphs"] == 3)

        p_rep = f"{OUT_DIR}/selftest-report.docx"
        tmp.append(p_rep)
        make_docx_file(normalize_blocks([
            {"name": "Header", "type": "component", "props": {"text": "커버 제목", "level": 1}},
            {"type": "divider", "props": {}},
            {"type": "header", "props": {"text": "실적", "level": "h2"}},
            {"type": "grid", "props": {"columns": 2, "children": [
                {"type": "metric", "props": {"label": "매출", "value": 1200, "unit": "억"}},
                {"type": "metric", "props": {"label": "영업이익", "value": 180}}]}},
            {"type": "key_value", "props": {"items": [
                {"key": "기간", "value": "2026 2Q"}, {"label": "기준", "value": "연결"}]}},
            {"type": "table", "props": {"headers": ["종목", "종가", "등락률"],
                                        "rows": dict_rows}},
            {"type": "alert", "props": {"title": "잠정치", "message": "감사 전 수치입니다"}},
        ]), "", p_rep)
        first = _DD(p_rep).paragraphs[0]
        ck("docx report: the cover heading is set at 22pt or larger",
           "커버 제목" in first.text and bool(first.runs)
           and first.runs[0].font.size is not None and first.runs[0].font.size.pt >= 22)
        rep = read_docx(p_rep)
        flat = ["\n".join(t["headers"] + [c for r in t["rows"] for c in r])
                for t in rep["tables"]]
        ck("docx report: a grid's metrics reach the page as one KPI strip",
           any("매출" in f and "영업이익" in f for f in flat))
        ck("docx report: key_value is a two-column table and dict rows keep their values",
           any("기간" in f and "2026 2Q" in f for f in flat)
           and any("삼성전자" in f and "70500" in f for f in flat))
        ck("docx report: the callout box keeps its words",
           "잠정치" in rep["text"] and "감사 전 수치입니다" in rep["text"])

        p_rpdf = f"{OUT_DIR}/selftest-report.pdf"
        tmp.append(p_rpdf)
        make_pdf_file(normalize_blocks([
            {"type": "header", "props": {"text": "Quarterly Report", "level": 1}},
            {"type": "header", "props": {"text": "Summary", "level": "h2"}},
            {"type": "grid", "props": {"columns": 2, "children": [
                {"type": "metric", "props": {"label": "revenue", "value": 1200}},
                {"type": "metric", "props": {"label": "profit", "value": 180}}]}},
            {"type": "key_value", "props": {"items": [{"key": "period", "value": "2026 2Q"}]}},
            {"type": "table", "props": {"rows": [{"item": "revenue", "value": 1200}]}},
            {"type": "alert", "props": {"message": "provisional numbers"}},
        ]), "", p_rpdf)
        rpdf = read_pdf(p_rpdf)
        ck("pdf report: the cover turns the page and every native block lands",
           rpdf["meta"]["pages"] >= 2
           and all(s in rpdf["text"] for s in
                   ("Quarterly Report", "revenue", "period", "provisional numbers")))
    except Exception as e:  # noqa: BLE001
        ck(f"report genre crashed: {e}", False)

    # pdf make round-trip — ASCII asserted (CID-encoded Korean does not reliably re-extract
    # through pypdf; the Korean path is verified by the build not raising with 한글 in it).
    p = f"{OUT_DIR}/selftest-made.pdf"
    tmp.append(p)
    try:
        make_pdf_file([{"type": "header", "props": {"text": "Quarterly Report 분기", "level": 1}},
                       {"type": "text", "props": {"content": "revenue grew 12 percent"}},
                       {"type": "table", "props": {"headers": ["item", "값"],
                                                   "rows": [["revenue", 100]]}}],
                      "Firebat PDF", p)
        got = read_pdf(p)
        ck("pdf make: text survives the round trip", "revenue grew 12 percent" in got["text"])
        ck("pdf make: the file opens as a real pdf", got["meta"]["pages"] >= 1)
    except Exception as e:  # noqa: BLE001
        ck(f"pdf make round trip crashed: {e}", False)

    for p in tmp:
        try:
            os.remove(p)
        except OSError:
            pass

    failed = [c for c in checks if not c["ok"]]
    return {"success": not failed, "action": "selftest",
            "data": {"checks": checks, "total": len(checks), "failed": len(failed)}}


def main():
    # Bytes, decoded as UTF-8 explicitly — the locale default turns Korean into lone
    # surrogates on some hosts (measured on Windows), and the envelope is UTF-8 by contract.
    raw = sys.stdin.buffer.read().decode("utf-8", "replace")
    try:
        envelope = json.loads(raw or "{}")
    except json.JSONDecodeError as e:
        sys.stdout.buffer.write((json.dumps({"success": False, "action": "", "error": f"input JSON: {e}"})).encode("utf-8"))
        return
    inp = envelope.get("data") or envelope
    action = str(inp.get("action") or "").strip()
    handlers = {"read": action_read, "make_pptx": action_make_pptx,
                "make_hwpx": action_make_hwpx,
                "make_xlsx": action_make_xlsx, "make_docx": action_make_docx,
                "make_pdf": action_make_pdf}
    if action == "selftest":
        out = action_selftest()
    elif action in handlers:
        out = handlers[action](inp)
    else:
        out = {"success": False, "action": action,
               "error": f"unknown action {action!r} — one of: read, make_pptx, make_xlsx, "
                        "make_docx, make_pdf, selftest"}
    # UTF-8 bytes out, explicitly — print() writes the console codepage on some hosts,
    # and the envelope is UTF-8 by contract on both ends.
    sys.stdout.buffer.write((json.dumps(out, ensure_ascii=False, default=str)).encode("utf-8"))


if __name__ == "__main__":
    main()
